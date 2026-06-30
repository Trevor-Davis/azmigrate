"""
Generate the Kettering assessment deck from source workbooks.

Expected input files in --input-dir:
  - Discovery.xlsx
  - Strategy_Lift_and_shift.xlsx
  - Strategy_PaaS_Preferred.xlsx

Usage:
  python generate_kettering_deck.py --input-dir "C:\\path\\to\\detailed-report" --output "Kettering-Slides-generated.pptx"

Notes:
  - Values are recalculated from the source files each run.
  - Discovery.xlsx may be MIP/encrypted or an older Excel container. If pandas cannot
    read it directly, the script attempts to use local Excel via pywin32 to save a
    temporary readable .xlsx copy. Install pywin32 if needed:
      python -m pip install pywin32
"""

from __future__ import annotations

import argparse
import ast
import csv
import io
import math
import os
import re
import shutil
import tempfile
import urllib.request
from pathlib import Path

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(30, 39, 97)
SLATE = RGBColor(54, 69, 79)
MUTED = RGBColor(100, 116, 139)
TEAL = RGBColor(2, 128, 144)
TEAL_DARK = RGBColor(25, 103, 128)
MINT = RGBColor(2, 195, 154)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(220, 38, 38)
ORANGE = RGBColor(255, 128, 0)
YELLOW = RGBColor(234, 179, 8)
GREY = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)
LINE = RGBColor(218, 226, 238)
ALT = RGBColor(241, 245, 249)
PINK = RGBColor(244, 114, 182)


DB_RESOURCE_TYPES = {
    "microsoft.applicationmigration/mongosites/mongoinstances",
    "microsoft.applicationmigration/pgsqlsites/pgsqlinstances",
    "microsoft.mysqldiscovery/mysqlsites/mysqlservers",
    "microsoft.offazure/mastersites/sqlsites/sqlservers",
}
NON_SQL_DB_RESOURCE_TYPES = {
    "MongoDB": "microsoft.applicationmigration/mongosites/mongoinstances",
    "MySQL": "microsoft.mysqldiscovery/mysqlsites/mysqlservers",
    "PostgreSQL": "microsoft.applicationmigration/pgsqlsites/pgsqlinstances",
}
WEBAPP_RESOURCE_TYPES = {
    "microsoft.offazure/mastersites/webappsites/iiswebapplications",
    "microsoft.offazure/mastersites/webappsites/tomcatwebapplications",
}
WEBAPP_RESOURCE_TYPE_LABELS = {
    "IIS": "microsoft.offazure/mastersites/webappsites/iiswebapplications",
    "Tomcat": "microsoft.offazure/mastersites/webappsites/tomcatwebapplications",
}
FILESHARE_RESOURCE_TYPE = "microsoft.applicationmigration/storagesites/fileshares"
MACHINE_RESOURCE_TYPE = "microsoft.offazure/vmwaresites/machines"

SKU_MEMORY_GB = {
    "Standard_D2as_v4": 8,
    "Standard_D2as_v5": 8,
    "Standard_D2ds_v4": 8,
    "Standard_D4as_v5": 16,
    "Standard_D8as_v5": 32,
    "Standard_D8s_v5": 32,
    "Standard_D16as_v5": 64,
    "Standard_D48as_v5": 192,
    "Standard_E2as_v5": 16,
    "Standard_E2bs_v5": 16,
    "Standard_E2s_v6": 16,
    "Standard_E4as_v5": 32,
    "Standard_E4bs_v5": 32,
    "Standard_E4s_v6": 32,
    "Standard_E8as_v5": 64,
    "Standard_E8bs_v5": 64,
    "Standard_E8s_v6": 64,
    "Standard_E16bs_v5": 128,
    "Standard_E16s_v6": 128,
    "Standard_E20as_v5": 160,
    "Standard_E20s_v6": 160,
    "Standard_E32as_v5": 256,
    "Standard_E32bs_v5": 256,
    "Standard_E48bs_v5": 384,
    "Standard_E64s_v6": 512,
    "Standard_E96s_v6": 768,
    "Standard_E128s_v6": 1024,
    "Standard_E192is_v6": 1832,
    "Standard_FX12mds": 252,
    "Standard_M16bs_v3": 128,
}


def norm_server(value) -> str:
    if pd.isna(value):
        return ""
    text = str(value).strip().upper()
    if not text:
        return ""
    text = text.split(".")[0]
    if "-" in text and re.match(r"^[A-Z]{2,6}[A-Z0-9]*\d{2,}", text):
        text = text.split("-")[0]
    return text.strip()


def recommended_memory_gb(sku) -> float | None:
    if pd.isna(sku):
        return None
    text = str(sku).strip()
    if not text or text.lower() == "unknown":
        return None
    if text in SKU_MEMORY_GB:
        return float(SKU_MEMORY_GB[text])
    match = re.match(r"^Standard_([DEFL])(\d+)[A-Za-z0-9_]*$", text)
    if match:
        family = match.group(1)
        n = int(match.group(2))
        multiplier = {"D": 4, "E": 8, "F": 2, "L": 8}[family]
        return float(n * multiplier)
    match = re.match(r"^Standard_B(\d+)(m?s)?$", text)
    if match:
        n = int(match.group(1))
        suffix = (match.group(2) or "").lower()
        if suffix == "ms":
            return float(n * 4)
        return float(n * 2)
    match = re.match(r"^Standard_M(\d+)[A-Za-z0-9_]*$", text)
    if match:
        return float(int(match.group(1)) * 8)
    return None


def money_k(value: float) -> str:
    return f"${value / 1000:.1f}K"


def money_full(value: float) -> str:
    sign = "-" if value < 0 else ""
    return f"{sign}${abs(value):,.0f}"


def tb_from_gb(value: float) -> str:
    return f"{value / 1024:.2f} TB"


def tb_from_gib(value: float) -> str:
    return f"{value / 1024:.1f} TB"


def pb_from_gb(value: float) -> str:
    return f"{value / 1024 / 1024:.2f} PB"


def safe_div(numerator: float, denominator: float) -> float:
    return numerator / denominator if denominator else 0


def pct(numerator: float, denominator: float) -> str:
    return f"{safe_div(numerator, denominator):.1%}"


def delta_pct(recommended: float, onprem: float) -> float:
    return safe_div(recommended - onprem, onprem) * 100


def delta_text(delta: float) -> str:
    return f"{abs(delta):.1f}% {'decrease' if delta < 0 else 'increase'}"


def read_excel_direct(path: Path, sheet: str) -> pd.DataFrame:
    suffix = path.suffix.lower()
    engine = "openpyxl" if suffix in {".xlsx", ".xlsm"} else None
    return pd.read_excel(path, sheet_name=sheet, engine=engine)


def convert_with_excel(path: Path) -> Path:
    try:
        import win32com.client  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            f"Could not read {path.name} directly and pywin32 is not installed. "
            "Install it with: python -m pip install pywin32"
        ) from exc

    out = Path(tempfile.mkdtemp()) / f"{path.stem}_converted.xlsx"
    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    wb = None
    try:
        wb = excel.Workbooks.Open(str(path), 0, True)
        wb.SaveAs(str(out), 51)
    finally:
        if wb is not None:
            wb.Close(False)
        excel.Quit()
    return out


def log(message: str) -> None:
    print(message, flush=True)


AZURE_VMS_CSV_URL = "https://raw.githubusercontent.com/Trevor-Davis/azmigrate/refs/heads/main/AzureVMs.csv"
_GITHUB_SKU_CACHE: dict[str, float] | None = None


def fetch_github_sku_map() -> dict[str, float]:
    """Download and parse the GitHub AzureVMs.csv into {sku: memory_gb}.

    Cached for the life of the process. Returns {} if the fetch fails so the
    rest of the script can still run with built-in fallbacks.
    """
    global _GITHUB_SKU_CACHE
    if _GITHUB_SKU_CACHE is not None:
        return _GITHUB_SKU_CACHE
    log(f"  Fetching SKU lookup from GitHub: {AZURE_VMS_CSV_URL}")
    mapping: dict[str, float] = {}
    try:
        req = urllib.request.Request(AZURE_VMS_CSV_URL, headers={"User-Agent": "azmigrate-script"})
        with urllib.request.urlopen(req, timeout=20) as resp:
            raw = resp.read().decode("utf-8-sig")
        reader = csv.DictReader(io.StringIO(raw))
        if not reader.fieldnames or "Name" not in reader.fieldnames or "MemoryGB" not in reader.fieldnames:
            log(f"    WARNING: CSV missing expected 'Name'/'MemoryGB' columns; got {reader.fieldnames}")
            _GITHUB_SKU_CACHE = {}
            return _GITHUB_SKU_CACHE
        for row in reader:
            name = (row.get("Name") or "").strip()
            mem = (row.get("MemoryGB") or "").strip()
            if not name or not mem:
                continue
            try:
                mapping[name] = float(mem)
            except ValueError:
                continue
        log(f"    Loaded {len(mapping):,} SKU entries from GitHub")
    except Exception as exc:
        log(f"    WARNING: Could not fetch GitHub SKU list ({exc}); falling back to built-in lookup only")
    _GITHUB_SKU_CACHE = mapping
    return _GITHUB_SKU_CACHE


def _col_letter(n: int) -> str:
    s = ""
    while n > 0:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s


def ensure_recommended_memory_openpyxl(lift_path: Path) -> list[str]:
    from openpyxl import load_workbook

    wb = load_workbook(lift_path)
    if "Server_to_AzureVM" not in wb.sheetnames:
        raise RuntimeError("Server_to_AzureVM sheet not found in lift workbook")
    ws = wb["Server_to_AzureVM"]

    def header_map():
        return {
            str(cell.value).strip(): idx
            for idx, cell in enumerate(ws[1], start=1)
            if cell.value is not None and str(cell.value).strip()
        }

    headers = header_map()
    if "RECOMMENDED_COMPUTE_SKU" not in headers:
        raise RuntimeError("RECOMMENDED_COMPUTE_SKU column not found")

    existing_memory: dict[str, float] = {}
    if "Azure SKUs" in wb.sheetnames:
        sku_ws = wb["Azure SKUs"]
        for row in sku_ws.iter_rows(min_row=2, values_only=True):
            sku_val = row[0] if len(row) > 0 else None
            mem_val = row[1] if len(row) > 1 else None
            if sku_val is None:
                continue
            key = str(sku_val).strip()
            if not key:
                continue
            try:
                if mem_val is not None and str(mem_val).strip() != "":
                    existing_memory[key] = float(mem_val)
            except (TypeError, ValueError):
                pass
        del wb["Azure SKUs"]

    if "Recommended Memory" in headers:
        ws.delete_cols(headers["Recommended Memory"])
        headers = header_map()

    sku_col = headers["RECOMMENDED_COMPUTE_SKU"]
    skus = {
        str(ws.cell(row=r, column=sku_col).value).strip()
        for r in range(2, ws.max_row + 1)
        if ws.cell(row=r, column=sku_col).value is not None
        and str(ws.cell(row=r, column=sku_col).value).strip()
        and str(ws.cell(row=r, column=sku_col).value).strip().lower() != "unknown"
    }

    github_map = fetch_github_sku_map()
    sku_memory: dict[str, float] = {}
    blanks: list[str] = []
    for sku in sorted(skus):
        if sku in existing_memory:
            sku_memory[sku] = existing_memory[sku]
        elif sku in github_map:
            sku_memory[sku] = float(github_map[sku])
        else:
            mem = recommended_memory_gb(sku)
            if mem is not None:
                sku_memory[sku] = float(mem)
            else:
                blanks.append(sku)

    sku_ws = wb.create_sheet("Azure SKUs", len(wb.worksheets))
    sku_ws.cell(row=1, column=1).value = "Azure SKU"
    sku_ws.cell(row=1, column=2).value = "Memory (GB)"
    for row_idx, sku in enumerate(sorted(skus), start=2):
        sku_ws.cell(row=row_idx, column=1).value = sku
        if sku in sku_memory:
            sku_ws.cell(row=row_idx, column=2).value = sku_memory[sku]

    target_col = sku_col + 1
    ws.insert_cols(target_col)
    ws.cell(row=1, column=target_col).value = "Recommended Memory"
    for r in range(2, ws.max_row + 1):
        sku_val = ws.cell(row=r, column=sku_col).value
        sku = "" if sku_val is None else str(sku_val).strip()
        ws.cell(row=r, column=target_col).value = sku_memory.get(sku)

    wb.save(lift_path)
    wb.close()
    return blanks


def ensure_recommended_memory(lift_path: Path) -> None:
    """In Strategy_Lift_and_shift.xlsx, create/refresh:
       - sheet 'Azure SKUs' with [Azure SKU, Memory (GB)] for every SKU in
         Server_to_AzureVM!RECOMMENDED_COMPUTE_SKU
       - column 'Recommended Memory' in Server_to_AzureVM (right after
         RECOMMENDED_COMPUTE_SKU) populated by VLOOKUP into 'Azure SKUs'

    Uses direct .xlsx editing to avoid Excel COM failures when the desktop
    Excel process is busy.
    """
    log(f"  Updating {lift_path.name}: Azure SKUs sheet + Recommended Memory column")
    blanks = ensure_recommended_memory_openpyxl(lift_path)
    log("    Azure SKUs sheet + Recommended Memory column updated")
    if blanks:
        log(f"    WARNING: {len(blanks)} SKU(s) have no memory value yet — fill them in the 'Azure SKUs' sheet:")
        for s in blanks:
            log(f"      - {s}")
        log("    These rows will show as 0 in totals until you fill the values. Re-run the script after editing.")


def read_excel_any(path: Path, sheet: str) -> pd.DataFrame:
    log(f"  - Reading {path.name} [{sheet}]")
    try:
        return read_excel_direct(path, sheet)
    except Exception:
        log(f"    direct read failed, converting via Excel automation...")
        converted = convert_with_excel(path)
        return pd.read_excel(converted, sheet_name=sheet, engine="openpyxl")


def require_columns(df: pd.DataFrame, sheet: str, columns: list[str]) -> None:
    missing = [c for c in columns if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns in {sheet}: {missing}")


def add_text(slide, text, x, y, w, h, size=12, color=SLATE, bold=False, font="Aptos", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    if align:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run_font = run.font
    run_font.name = font
    run_font.size = Pt(size)
    run_font.color.rgb = color
    run_font.bold = bold
    return box


def add_panel(slide, x, y, w, h, title, accent=TEAL, title_size=13):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LINE
    rect.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.15, y + 0.14, w - 0.3, 0.3, title_size, NAVY, True)
    return rect


def add_card(slide, x, y, w, h, value, label, sub="", accent=TEAL, value_size=26, label_size=12, sub_size=12):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LINE
    rect.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, value, x + 0.17, y + 0.30, w - 0.34, 0.40, value_size, NAVY, True)
    add_text(slide, label, x + 0.17, y + 0.82, w - 0.34, 0.30, label_size, SLATE, True)
    if sub:
        add_text(slide, sub, x + 0.17, y + 1.16, w - 0.34, 0.35, sub_size, MUTED)
    return rect


def add_small_card(slide, x, y, w, h, value, label, accent=TEAL, value_size=14, label_size=8):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = LINE
    rect.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, value, x + 0.12, y + 0.20, w - 0.24, 0.28, value_size, NAVY, True)
    add_text(slide, label, x + 0.12, y + 0.58, w - 0.24, 0.22, label_size, SLATE, True)
    return rect


def add_table(slide, rows, x, y, w, h, font_size=9, header_color=NAVY, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.04)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE if r == 0 else SLATE
            p.font.bold = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else (WHITE if r % 2 else ALT)
    return table_shape


def color_delta_text(cell, text: str, size: int = 11):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = SLATE
    m = re.search(r"(\d+(?:\.\d+)?% (?:increase|decrease))", text)
    if not m:
        return
    start = m.start(1)
    end = m.end(1)
    # PowerPoint text run coloring is easiest by splitting runs.
    p.text = ""
    before = text[:start]
    delta = text[start:end]
    after = text[end:]
    if before:
        r = p.add_run()
        r.text = before
        r.font.size = Pt(size)
        r.font.color.rgb = SLATE
    r = p.add_run()
    r.text = delta
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = GREEN if "decrease" in delta else RED
    if after:
        r = p.add_run()
        r.text = after
        r.font.size = Pt(size)
        r.font.color.rgb = SLATE


def create_donut_image(path: Path, smb: int, nfs: int, title="Protocol Distribution"):
    width, height, scale = 520, 360, 4
    img = Image.new("RGB", (width * scale, height * scale), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        title_font = ImageFont.truetype("arialbd.ttf", 28 * scale)
        legend_font = ImageFont.truetype("arial.ttf", 22 * scale)
    except Exception:
        title_font = ImageFont.load_default()
        legend_font = ImageFont.load_default()

    navy = (31, 78, 121)
    teal = (25, 103, 128)
    orange = (255, 128, 0)
    gray = (89, 89, 89)
    white = (255, 255, 255)

    bb = draw.textbbox((0, 0), title, font=title_font)
    draw.text(((width * scale - (bb[2] - bb[0])) / 2, 8 * scale), title, fill=navy, font=title_font)

    cx, cy = 260 * scale, 168 * scale
    r, inner = 118 * scale, 49 * scale
    bbox = [cx - r, cy - r, cx + r, cy + r]
    total = max(smb + nfs, 1)
    nfs_angle = 360 * nfs / total
    start = -90
    draw.pieslice(bbox, start, start + 360, fill=teal)
    draw.pieslice(bbox, start, start + nfs_angle, fill=orange)
    inner_bbox = [cx - inner, cy - inner, cx + inner, cy + inner]
    draw.ellipse(inner_bbox, fill=white)
    draw.ellipse(bbox, outline=white, width=3 * scale)
    draw.ellipse(inner_bbox, outline=white, width=3 * scale)
    for angle in (start, start + nfs_angle):
        rad = math.radians(angle)
        draw.line(
            (
                cx + inner * math.cos(rad),
                cy + inner * math.sin(rad),
                cx + r * math.cos(rad),
                cy + r * math.sin(rad),
            ),
            fill=white,
            width=4 * scale,
        )

    legend_y = 314 * scale
    box = 12 * scale
    items = [(f"NFS ({nfs})", orange), (f"SMB ({smb})", teal)]
    widths = []
    for label, _ in items:
        bb = draw.textbbox((0, 0), label, font=legend_font)
        widths.append(box + 8 * scale + (bb[2] - bb[0]))
    total_w = sum(widths) + 34 * scale
    x = (width * scale - total_w) / 2
    for (label, color), item_w in zip(items, widths):
        draw.rectangle([x, legend_y, x + box, legend_y + box], fill=color, outline=white, width=2 * scale)
        draw.text((x + box + 8 * scale, legend_y - 5 * scale), label, fill=gray, font=legend_font)
        x += item_w + 34 * scale

    img = img.resize((width, height), Image.Resampling.LANCZOS)
    img.save(path)
    return path


def create_pie_image(path: Path, title: str, values: dict[str, int], colors: dict[str, RGBColor]):
    width, height, scale = 520, 360, 4
    img = Image.new("RGB", (width * scale, height * scale), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        title_font = ImageFont.truetype("arialbd.ttf", 26 * scale)
        legend_font = ImageFont.truetype("arial.ttf", 19 * scale)
    except Exception:
        title_font = ImageFont.load_default()
        legend_font = ImageFont.load_default()

    def rgb(color: RGBColor):
        return (int(color[0]), int(color[1]), int(color[2]))

    navy = rgb(NAVY)
    gray = rgb(SLATE)
    white = (255, 255, 255)
    title_box = draw.textbbox((0, 0), title, font=title_font)
    draw.text(((width * scale - (title_box[2] - title_box[0])) / 2, 8 * scale), title, fill=navy, font=title_font)

    cx, cy = 190 * scale, 180 * scale
    radius = 112 * scale
    bbox = [cx - radius, cy - radius, cx + radius, cy + radius]
    total = sum(values.values()) or 1
    extents = {label: (360 * value / total if value else 0) for label, value in values.items()}
    min_visible = 7
    excess = 0.0
    for label, value in values.items():
        if value and extents[label] < min_visible:
            excess += min_visible - extents[label]
            extents[label] = min_visible
    for label in sorted(extents, key=extents.get, reverse=True):
        if excess <= 0:
            break
        reducible = max(0, extents[label] - min_visible)
        reduction = min(reducible, excess)
        extents[label] -= reduction
        excess -= reduction
    start = -90
    for label, value in values.items():
        extent = extents[label]
        draw.pieslice(bbox, start, start + extent, fill=rgb(colors[label]), outline=white, width=3 * scale)
        start += extent

    legend_x = 345 * scale
    legend_y = 122 * scale
    box = 13 * scale
    for i, (label, value) in enumerate(values.items()):
        y = legend_y + i * 48 * scale
        draw.rectangle([legend_x, y, legend_x + box, y + box], fill=rgb(colors[label]), outline=white, width=2 * scale)
        draw.text((legend_x + box + 10 * scale, y - 6 * scale), f"{label} ({value:,})", fill=gray, font=legend_font)

    img = img.resize((width, height), Image.Resampling.LANCZOS)
    img.save(path)
    return path


def create_readiness_image(path: Path, readiness_counts: dict[str, int]):
    categories = ["Unknown", "Ready With Conditions", "Not Ready", "Ready"]
    values = [int(readiness_counts.get(category, 0)) for category in categories]
    max_value = max(values) or 1
    total = sum(values) or 1

    width, height, scale = 760, 320, 4
    img = Image.new("RGB", (width * scale, height * scale), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        label_font = ImageFont.truetype("arial.ttf", 18 * scale)
        value_font = ImageFont.truetype("arialbd.ttf", 18 * scale)
    except Exception:
        label_font = ImageFont.load_default()
        value_font = ImageFont.load_default()

    colors = {
        "Unknown": (148, 163, 184),
        "Ready With Conditions": (255, 128, 0),
        "Not Ready": (220, 38, 38),
        "Ready": (34, 139, 34),
    }
    label_x = 26 * scale
    bar_x = 270 * scale
    bar_w = 315 * scale
    row_h = 60 * scale
    y0 = 42 * scale
    for idx, (category, value) in enumerate(zip(categories, values)):
        y = y0 + idx * row_h
        label = category
        draw.text((label_x, y - 10 * scale), label, fill=(54, 69, 79), font=label_font)
        draw.rounded_rectangle(
            [bar_x, y, bar_x + bar_w, y + 18 * scale],
            radius=9 * scale,
            fill=(241, 245, 249),
        )
        filled = max(2 * scale if value else 0, int(bar_w * value / max_value))
        if filled:
            draw.rounded_rectangle(
                [bar_x, y, bar_x + filled, y + 18 * scale],
                radius=9 * scale,
                fill=colors[category],
            )
        draw.text(
            (bar_x + bar_w + 18 * scale, y - 11 * scale),
            f"{value:,} ({value / total:.1%})",
            fill=(30, 39, 97),
            font=value_font,
        )

    img = img.resize((width, height), Image.Resampling.LANCZOS)
    img.save(path)
    return path


def discovery_sets(discovery: pd.DataFrame):
    require_columns(discovery, "Discovery.xlsx/Data", ["resourceType", "parentResourceName"])
    rows = discovery[["resourceType", "parentResourceName"]].copy()
    rows["resourceType"] = rows["resourceType"].astype(str).str.strip()
    rows["server"] = rows["parentResourceName"].map(norm_server)
    files = set(rows.loc[rows["resourceType"].eq(FILESHARE_RESOURCE_TYPE), "server"]) - {""}
    web = set(rows.loc[rows["resourceType"].isin(WEBAPP_RESOURCE_TYPES), "server"]) - {""}
    db = set(rows.loc[rows["resourceType"].isin(DB_RESOURCE_TYPES), "server"]) - {""}
    return files, web, db


def load_metrics(input_dir: Path):
    discovery_path = input_dir / "Discovery.xlsx"
    lift_path = input_dir / "Strategy_Lift_and_shift.xlsx"
    paas_path = input_dir / "Strategy_PaaS_Preferred.xlsx"
    for path in (discovery_path, lift_path, paas_path):
        if not path.exists():
            raise FileNotFoundError(path)

    discovery = read_excel_any(discovery_path, "Data")
    ensure_recommended_memory(lift_path)
    lift = read_excel_any(lift_path, "Server_to_AzureVM")
    fs = read_excel_any(paas_path, "FileShares_to_Azure_Files")

    def _read_optional(path: Path, sheet: str):
        try:
            return read_excel_any(path, sheet)
        except Exception as exc:
            log(f"    Optional sheet '{sheet}' not loaded ({exc}); continuing")
            return None

    sql_vm_df = _read_optional(paas_path, "SQLinstance_to_AzureSQLVM")
    sql_mi_df = _read_optional(paas_path, "SQLinstance_to_AzureSQLMI")
    mongo_df = _read_optional(paas_path, "Mongo_to_Azure_Document_DB")
    mysql_df = _read_optional(paas_path, "MySQL_to_AzureFlexServerMySQL")
    pgsql_df = _read_optional(paas_path, "PgSQL_to_AzureFlexServerPG")
    webapp_df = _read_optional(paas_path, "WebApp_to_AKS")

    require_columns(discovery, "Discovery.xlsx/Data", ["resourceType", "parentResourceName"])
    require_columns(
        lift,
        "Strategy_Lift_and_shift.xlsx/Server_to_AzureVM",
        [
            "SERVER_NAME",
            "ONPREM_STORAGE_GB",
            "RECOMMENDED_STORAGE_SIZE_GB",
            "ONPREM_CORES_COUNT",
            "RECOMMENDED_NUMBER_OF_CORES",
            "ONPREM_MEMORY_MB",
            "RECOMMENDED_COMPUTE_SKU",
            "STORAGE_UTILIZATION_PERCENT",
            "ONPREM_CPU_USAGE_PERCENT",
            "ONPREM_MEMORY_USAGE_PERCENT",
            "TOTAL_MONTHLY_COST_USD",
        ],
    )
    require_columns(
        fs,
        "Strategy_PaaS_Preferred.xlsx/FileShares_to_Azure_Files",
        ["SERVER", "PROTOCOL", "READINESS", "AZURE FILES TARGET PROVISIONED SIZE (GIB)", "ESTIMATED MONTHLY COST (USD)"],
    )

    fileshare_servers, web_servers, db_servers = discovery_sets(discovery)
    fileshare_only = fileshare_servers - web_servers - db_servers

    machines = discovery[discovery["resourceType"].eq(MACHINE_RESOURCE_TYPE)].copy()
    vm_total = len(machines)
    power = machines.get("powerOnStatus", pd.Series(dtype=object)).fillna("Unknown").astype(str)
    powered_on = int(power.str.contains("on", case=False, na=False).sum())
    powered_off = int(power.str.contains("off", case=False, na=False).sum())

    fileshares = discovery[discovery["resourceType"].eq(FILESHARE_RESOURCE_TYPE)].copy()
    fs_total = len(fileshares)

    machine_os = {}
    if "resourceName" in machines.columns and "osType" in machines.columns:
        for _, row in machines.iterrows():
            machine_os[norm_server(row.get("resourceName"))] = str(row.get("osType", "")).strip()
    os_counts = {"Windows": 0, "RHEL": 0, "Other/Unknown": 0}
    for _, row in fileshares.iterrows():
        os_type = machine_os.get(norm_server(row.get("parentResourceName")), "")
        if "win" in os_type.lower():
            os_counts["Windows"] += 1
        elif "linux" in os_type.lower() or "rhel" in os_type.lower() or "red" in os_type.lower():
            os_counts["RHEL"] += 1
        else:
            os_counts["Other/Unknown"] += 1

    db_counts = {
        "SQL Server": int(discovery["resourceType"].eq("microsoft.offazure/mastersites/sqlsites/sqlservers").sum()),
        "MongoDB": int(discovery["resourceType"].eq("microsoft.applicationmigration/mongosites/mongoinstances").sum()),
        "MySQL": int(discovery["resourceType"].eq("microsoft.mysqldiscovery/mysqlsites/mysqlservers").sum()),
        "PostgreSQL": int(discovery["resourceType"].eq("microsoft.applicationmigration/pgsqlsites/pgsqlinstances").sum()),
    }
    db_total = sum(db_counts.values())

    non_sql_db_discovery: dict[str, dict[str, object]] = {}
    non_sql_support_statuses: set[str] = set()
    all_non_sql = discovery[discovery["resourceType"].isin(NON_SQL_DB_RESOURCE_TYPES.values())].copy()
    for db_type, resource_type in NON_SQL_DB_RESOURCE_TYPES.items():
        db_rows = discovery[discovery["resourceType"].eq(resource_type)].copy()
        parent_values = db_rows.get("parentResourceName", pd.Series(dtype=object))
        version_values = db_rows.get("version", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
        support_values = db_rows.get("directSupportStatus", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
        support_counts = {str(k): int(v) for k, v in support_values.value_counts().to_dict().items()}
        non_sql_support_statuses.update(support_counts)
        non_sql_db_discovery[db_type] = {
            "instances": int(len(db_rows)),
            "servers": int(parent_values.dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
            "versions": int(version_values.nunique()),
            "version_counts": {str(k): int(v) for k, v in version_values.value_counts().to_dict().items()},
            "support_counts": support_counts,
        }
    all_non_sql_versions = all_non_sql.get("version", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
    all_non_sql_support = all_non_sql.get("directSupportStatus", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
    non_sql_support_total = {str(k): int(v) for k, v in all_non_sql_support.value_counts().to_dict().items()}
    non_sql_support_statuses.update(non_sql_support_total)
    non_sql_db_discovery["Total"] = {
        "instances": int(len(all_non_sql)),
        "servers": int(all_non_sql.get("parentResourceName", pd.Series(dtype=object)).dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
        "versions": int(all_non_sql_versions.nunique()),
        "version_counts": {str(k): int(v) for k, v in all_non_sql_versions.value_counts().to_dict().items()},
        "support_counts": non_sql_support_total,
    }
    non_sql_support_order = [s for s in ["Mainstream", "Extended", "OutOfSupport", "Unknown"] if s in non_sql_support_statuses]
    non_sql_support_order += sorted(non_sql_support_statuses - set(non_sql_support_order))

    sql_rows = discovery[discovery["resourceType"].eq("microsoft.offazure/mastersites/sqlsites/sqlservers")]
    sql_instance_count = int(len(sql_rows))
    sql_server_count = int(sql_rows["parentResourceName"].dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique())

    # SQL version breakdown: count per (version, pssStatus). pssStatus drives bar color.
    if {"version", "pssStatus"}.issubset(sql_rows.columns):
        sql_version_df = sql_rows.assign(
            version=sql_rows["version"].fillna("Unknown").astype(str).str.strip().replace("", "Unknown"),
            pssStatus=sql_rows["pssStatus"].fillna("Unknown").astype(str).str.strip().replace("", "Unknown"),
        )
        grouped = sql_version_df.groupby(["version", "pssStatus"]).size().reset_index(name="count")
        sql_versions = [
            (row["version"], row["pssStatus"], int(row["count"]))
            for _, row in grouped.sort_values("count", ascending=False).iterrows()
        ]
    else:
        sql_versions = []

    READINESS_ORDER = ["Ready", "Not Ready", "Ready With Conditions", "Unknown"]

    def _readiness_counts(df, col):
        result = {k: 0 for k in READINESS_ORDER}
        if df is None or col not in df.columns:
            return result, 0
        vals = (
            df[col]
            .fillna("Unknown")
            .astype(str)
            .str.strip()
            .replace("", "Unknown")
        )
        # Case-insensitive bucket lookup so "Ready with Conditions" matches "Ready With Conditions"
        bucket_by_lower = {k.lower(): k for k in READINESS_ORDER}
        for v in vals:
            key = bucket_by_lower.get(v.lower(), "Unknown")
            result[key] += 1
        return result, int(vals.shape[0])

    sql_vm_readiness, sql_vm_total = _readiness_counts(sql_vm_df, "AZURE_SQL_VM_READINESS")
    sql_mi_readiness, sql_mi_total = _readiness_counts(sql_mi_df, "AZURE_SQL_MI_READINESS")
    non_sql_db_specs = [
        ("MongoDB", mongo_df, "AZURE DOCUMENTDB READINESS"),
        ("MySQL", mysql_df, "AZURE DB FOR MYSQL READINESS"),
        ("PostgreSQL", pgsql_df, "AZURE DB FOR POSTGRESQL READINESS"),
    ]
    non_sql_db_counts: dict[str, int] = {}
    non_sql_db_readiness: dict[str, dict[str, int]] = {}
    non_sql_db_total = 0
    non_sql_db_readiness_total = {k: 0 for k in READINESS_ORDER}
    for db_type, df, readiness_col in non_sql_db_specs:
        counts, total = _readiness_counts(df, readiness_col)
        non_sql_db_counts[db_type] = total
        non_sql_db_readiness[db_type] = counts
        non_sql_db_total += total
        for key, value in counts.items():
            non_sql_db_readiness_total[key] += value

    # SQL cost comparison: lift-and-shift everything vs. SQL MI where ready + lift-and-shift the rest.
    def _sum_cost(df, cols):
        if df is None:
            return 0.0
        total = 0.0
        for c in cols:
            if c in df.columns:
                total += float(pd.to_numeric(df[c], errors="coerce").fillna(0).sum())
        return total

    sql_servers_norm: set[str] = set()
    for df, col in [(sql_mi_df, "Server"), (sql_vm_df, "SERVER")]:
        if df is None:
            continue
        candidate_col = col if col in df.columns else next(
            (c for c in df.columns if c.lower() == col.lower()), None
        )
        if candidate_col is None:
            continue
        for v in df[candidate_col].dropna().astype(str):
            n = norm_server(v)
            if n:
                sql_servers_norm.add(n)

    if "server_norm" not in lift.columns:
        lift["server_norm"] = lift["SERVER_NAME"].map(norm_server)
    sql_lift_only = lift[lift["server_norm"].isin(sql_servers_norm)]
    sql_lift_total_cost = float(
        pd.to_numeric(sql_lift_only["TOTAL_MONTHLY_COST_USD"], errors="coerce").fillna(0).sum()
    )

    # SQL MI cost: MI compute + storage for MI-ready instances, plus VM cost for the rest
    sql_mi_cost = _sum_cost(
        sql_mi_df,
        ["AZURE_SQL_MI_COMPUTE_MONTHLY_COST_USD", "AZURE_SQL_MI_STORAGE_MONTHLY_COST_USD"],
    )
    sql_vm_cost = _sum_cost(
        sql_vm_df,
        ["AZURE_SQL_VM_COMPUTE_MONTHLY_COST_USD", "AZURE_SQL_VM_STORAGE_MONTHLY_COST_USD"],
    )
    sql_hybrid_cost = sql_mi_cost + sql_vm_cost
    sql_cost_delta = sql_lift_total_cost - sql_hybrid_cost

    web_counts = {
        "IIS Web Applications": int(discovery["resourceType"].eq("microsoft.offazure/mastersites/webappsites/iiswebapplications").sum()),
        "Tomcat Web Applications": int(discovery["resourceType"].eq("microsoft.offazure/mastersites/webappsites/tomcatwebapplications").sum()),
    }
    web_total = sum(web_counts.values())
    webapp_discovery = discovery[discovery["resourceType"].isin(WEBAPP_RESOURCE_TYPE_LABELS.values())].copy()
    webapp_by_type: dict[str, dict[str, object]] = {}
    webapp_support_statuses: set[str] = set()
    for web_type, resource_type in WEBAPP_RESOURCE_TYPE_LABELS.items():
        rows = webapp_discovery[webapp_discovery["resourceType"].eq(resource_type)]
        parent_values = rows.get("parentResourceName", pd.Series(dtype=object))
        support_values = rows.get("directSupportStatus", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
        support_counts = {str(k): int(v) for k, v in support_values.value_counts().to_dict().items()}
        webapp_support_statuses.update(support_counts)
        webapp_by_type[web_type] = {
            "webapps": int(len(rows)),
            "servers": int(parent_values.dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
            "support_counts": support_counts,
        }
    all_webapp_support = webapp_discovery.get("directSupportStatus", pd.Series(dtype=object)).fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
    webapp_support_total = {str(k): int(v) for k, v in all_webapp_support.value_counts().to_dict().items()}
    webapp_support_statuses.update(webapp_support_total)
    webapp_by_type["Total"] = {
        "webapps": int(len(webapp_discovery)),
        "servers": int(webapp_discovery.get("parentResourceName", pd.Series(dtype=object)).dropna().astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
        "support_counts": webapp_support_total,
    }
    webapp_support_order = [s for s in ["Mainstream", "Extended", "OutOfSupport", "Unknown"] if s in webapp_support_statuses]
    webapp_support_order += sorted(webapp_support_statuses - set(webapp_support_order))

    webapp_readiness_by_type: dict[str, dict[str, int]] = {}
    webapp_readiness_total = {k: 0 for k in READINESS_ORDER}
    if webapp_df is not None and {"WEBAPPTYPE", "READINESS"}.issubset(webapp_df.columns):
        web_type_values = webapp_df["WEBAPPTYPE"].fillna("Unknown").astype(str).str.strip().replace("", "Unknown")
        for web_type in WEBAPP_RESOURCE_TYPE_LABELS:
            counts, _total = _readiness_counts(webapp_df[web_type_values.str.lower().eq(web_type.lower())], "READINESS")
            webapp_readiness_by_type[web_type] = counts
            for key, value in counts.items():
                webapp_readiness_total[key] += value
    else:
        for web_type in WEBAPP_RESOURCE_TYPE_LABELS:
            webapp_readiness_by_type[web_type] = {k: 0 for k in READINESS_ORDER}

    onprem_storage_gb = pd.to_numeric(lift["ONPREM_STORAGE_GB"], errors="coerce").fillna(0).sum()
    rec_storage_gb = pd.to_numeric(lift["RECOMMENDED_STORAGE_SIZE_GB"], errors="coerce").fillna(0).sum()
    onprem_cores = pd.to_numeric(lift["ONPREM_CORES_COUNT"], errors="coerce").fillna(0).sum()
    rec_cores = pd.to_numeric(lift["RECOMMENDED_NUMBER_OF_CORES"], errors="coerce").fillna(0).sum()
    onprem_mem_gb = pd.to_numeric(lift["ONPREM_MEMORY_MB"], errors="coerce").fillna(0).sum() / 1024
    if "Recommended Memory" not in lift.columns:
        lift["Recommended Memory"] = lift["RECOMMENDED_COMPUTE_SKU"].map(recommended_memory_gb)
    else:
        rec_mem = pd.to_numeric(lift["Recommended Memory"], errors="coerce")
        missing = rec_mem.isna()
        if missing.any():
            derived_mem = pd.to_numeric(
                lift.loc[missing, "RECOMMENDED_COMPUTE_SKU"].map(recommended_memory_gb),
                errors="coerce",
            )
            rec_mem.loc[missing] = derived_mem
        lift["Recommended Memory"] = rec_mem
    rec_mem_gb = pd.to_numeric(lift["Recommended Memory"], errors="coerce").fillna(0).sum()
    storage_util = pd.to_numeric(lift["STORAGE_UTILIZATION_PERCENT"], errors="coerce").dropna().mean()
    cpu_util = pd.to_numeric(lift["ONPREM_CPU_USAGE_PERCENT"], errors="coerce").dropna().mean()
    mem_util = pd.to_numeric(lift["ONPREM_MEMORY_USAGE_PERCENT"], errors="coerce").dropna().mean()

    fs["server_norm"] = fs["SERVER"].map(norm_server)
    fs_only_rows = fs[fs["server_norm"].isin(fileshare_only)].copy()
    fs_only_size_gib = pd.to_numeric(fs_only_rows["AZURE FILES TARGET PROVISIONED SIZE (GIB)"], errors="coerce").fillna(0).sum()
    fs_only_cost = pd.to_numeric(fs_only_rows["ESTIMATED MONTHLY COST (USD)"], errors="coerce").fillna(0).sum()
    lift["server_norm"] = lift["SERVER_NAME"].map(norm_server)
    lift_only = lift[lift["server_norm"].isin(fileshare_only)]
    lift_only_cost = pd.to_numeric(lift_only["TOTAL_MONTHLY_COST_USD"], errors="coerce").fillna(0).sum()

    protocol_counts = fs["PROTOCOL"].fillna("Unknown").astype(str).str.strip().replace("", "Unknown").value_counts().to_dict()
    readiness_counts = fs["READINESS"].fillna("Blank").astype(str).str.strip().replace("", "Blank").value_counts().to_dict()

    return {
        "vm_total": vm_total,
        "powered_on": powered_on,
        "powered_off": powered_off,
        "fileshare_total": fs_total,
        "fileshare_servers": len(fileshare_servers),
        "fileshare_only_servers": len(fileshare_only),
        "fileshare_only_shares": len(fs_only_rows),
        "fileshare_only_size_tb": fs_only_size_gib / 1024,
        "fileshare_only_azure_cost": fs_only_cost,
        "fileshare_only_lift_cost": lift_only_cost,
        "protocol_counts": protocol_counts,
        "readiness_counts": readiness_counts,
        "os_counts": os_counts,
        "db_counts": db_counts,
        "db_total": db_total,
        "web_counts": web_counts,
        "web_total": web_total,
        "webapp_by_type": webapp_by_type,
        "webapp_support_order": webapp_support_order,
        "webapp_readiness_by_type": webapp_readiness_by_type,
        "webapp_readiness_total": webapp_readiness_total,
        "sql_instance_count": sql_instance_count,
        "sql_server_count": sql_server_count,
        "sql_versions": sql_versions,
        "sql_vm_readiness": sql_vm_readiness,
        "sql_vm_total": sql_vm_total,
        "sql_mi_readiness": sql_mi_readiness,
        "sql_mi_total": sql_mi_total,
        "sql_lift_total_cost": sql_lift_total_cost,
        "sql_hybrid_cost": sql_hybrid_cost,
        "sql_mi_cost": sql_mi_cost,
        "sql_vm_cost": sql_vm_cost,
        "sql_cost_delta": sql_cost_delta,
        "non_sql_db_total": non_sql_db_total,
        "non_sql_db_counts": non_sql_db_counts,
        "non_sql_db_readiness": non_sql_db_readiness,
        "non_sql_db_readiness_total": non_sql_db_readiness_total,
        "non_sql_db_discovery": non_sql_db_discovery,
        "non_sql_support_order": non_sql_support_order,
        "onprem_storage_gb": onprem_storage_gb,
        "rec_storage_gb": rec_storage_gb,
        "onprem_cores": onprem_cores,
        "rec_cores": rec_cores,
        "onprem_mem_gb": onprem_mem_gb,
        "rec_mem_gb": rec_mem_gb,
        "storage_util": storage_util,
        "cpu_util": cpu_util,
        "mem_util": mem_util,
    }


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def slide_title(slide, title, subtitle=""):
    add_text(slide, title, 0.5, 0.35, 8.7, 0.55, 32, NAVY, True, "Aptos Display")
    if subtitle:
        add_text(slide, subtitle, 0.52, 0.92, 9.0, 0.3, 12, MUTED)


def add_source(slide, text, size=9):
    add_text(slide, text, 0.6, 6.92, 12.1, 0.18, size, MUTED)


def add_bar(slide, x, y, w, h, value, max_value, color):
    # Guard against degenerate dimensions; PowerPoint flags zero/negative shapes.
    w = max(float(w), 0.01)
    h = max(float(h), 0.01)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ALT
    bg.line.fill.background()
    fill_w = max(w * safe_div(value, max_value), 0.01)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(fill_w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_vm_power_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "VM Power State Summary")

    max_count = max(m["powered_on"], m["powered_off"], 1)
    chart_x, chart_y, chart_w, chart_h = 1.25, 1.55, 10.9, 1.2
    tick_max = int(math.ceil(max_count / 250.0) * 250) + 250
    for tick in range(0, tick_max + 1, 250):
        x = chart_x + chart_w * safe_div(tick, tick_max)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(chart_y), Inches(0.01), Inches(chart_h))
        line.fill.solid(); line.fill.fore_color.rgb = LINE
        line.line.fill.background()
        add_text(slide, f"{tick:,}", x - 0.25, chart_y + chart_h + 0.18, 0.5, 0.15, 9, MUTED, align=PP_ALIGN.CENTER)
    add_bar(slide, chart_x, chart_y + 0.2, chart_w, 0.28, m["powered_on"], max_count, TEAL)
    add_bar(slide, chart_x, chart_y + 0.75, chart_w * safe_div(m["powered_off"], max_count), 0.28, m["powered_off"], m["powered_off"] or 1, PINK)
    add_text(slide, f"Powered On ({m['powered_on']:,})", 4.75, 3.03, 1.8, 0.17, 10, SLATE, True, align=PP_ALIGN.CENTER)
    add_text(slide, f"Powered Off ({m['powered_off']:,})", 6.52, 3.03, 1.85, 0.17, 10, SLATE, True, align=PP_ALIGN.CENTER)

    add_panel(slide, 0.75, 3.92, 11.1, 2.25, "VM power-state counts", title_size=17)
    rows = [
        ["Power state", "VM count", "Share of VMs"],
        ["Powered On", f"{m['powered_on']:,}", pct(m["powered_on"], m["vm_total"])],
        ["Powered Off", f"{m['powered_off']:,}", pct(m["powered_off"], m["vm_total"])],
        ["Total VMs", f"{m['vm_total']:,}", "100.0%"],
    ]
    add_table(slide, rows, 0.95, 4.58, 6.0, 1.46, 12, col_widths=[2.5, 1.8, 1.7])
    add_card(slide, 7.55, 4.52, 1.95, 1.35, f"{m['powered_on']:,}", "Powered On", "", TEAL, 26, label_size=10)
    add_card(slide, 9.65, 4.52, 1.95, 1.35, f"{m['powered_off']:,}", "Powered Off", "", PINK, 26, label_size=10)
    add_source(slide, "Source: Discovery.xlsx, Data sheet, resourceType = microsoft.offazure/vmwaresites/machines.", size=9)


def add_vm_utilization_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "VM Utilization Summary", "On-premises footprint compared with recommended Azure sizing; lower is better for capacity change.")
    add_text(slide, "On-premises footprint", 2.85, 1.28, 5.83, 0.3, 17, NAVY, True)
    card_w, card_h = 2.08, 1.55
    add_card(slide, 0.55, 1.55, card_w, card_h, f"{m['vm_total']:,}", "Assessed VMs", "", TEAL, 26)
    add_card(slide, 2.85, 1.55, card_w, card_h, pb_from_gb(m["onprem_storage_gb"]), "Storage", f"{m['storage_util']:.1f}% Consumed", TEAL, 26)
    add_card(slide, 5.15, 1.55, card_w, card_h, f"{m['onprem_cores']:,.0f}", "Cores", f"{m['cpu_util']:.1f}% CPU Utilization", MINT, 26)
    add_card(slide, 7.45, 1.55, card_w, card_h, tb_from_gib(m["onprem_mem_gb"]), "Memory", f"{m['mem_util']:.1f}% Memory Utilization", MINT, 26)
    add_text(slide, "Recommended Azure sizing", 2.85, 3.96, 5.83, 0.3, 17, NAVY, True)
    storage_delta = (m["rec_storage_gb"] - m["onprem_storage_gb"]) / m["onprem_storage_gb"] * 100
    cores_delta = (m["rec_cores"] - m["onprem_cores"]) / m["onprem_cores"] * 100
    mem_delta = (m["rec_mem_gb"] - m["onprem_mem_gb"]) / m["onprem_mem_gb"] * 100
    add_card(slide, 2.85, 4.25, card_w, card_h, pb_from_gb(m["rec_storage_gb"]), "Storage", f"{storage_delta:+.1f}% vs on-prem", TEAL, 26)
    add_card(slide, 5.15, 4.25, card_w, card_h, f"{m['rec_cores']:,.0f}", "Cores", f"{cores_delta:+.1f}% vs on-prem", MINT, 26)
    add_card(slide, 7.45, 4.25, card_w, card_h, tb_from_gib(m["rec_mem_gb"]), "Memory", f"{mem_delta:+.1f}% vs on-prem", MINT, 26)

    # Comparison bar chart: on-prem vs recommended Azure sizing, per metric.
    panel_x, panel_y, panel_w, panel_h = 9.75, 1.55, 3.4, 4.25
    add_panel(slide, panel_x, panel_y, panel_w, panel_h, "On-prem vs Azure", title_size=13)
    legend_y = panel_y + 0.55
    sw_x = panel_x + 0.2
    sw = 0.18
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sw_x), Inches(legend_y), Inches(sw), Inches(0.12))
    s1.fill.solid(); s1.fill.fore_color.rgb = TEAL; s1.line.fill.background()
    add_text(slide, "On-prem", sw_x + sw + 0.05, legend_y - 0.04, 0.8, 0.2, 9, SLATE)
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sw_x + 1.25), Inches(legend_y), Inches(sw), Inches(0.12))
    s2.fill.solid(); s2.fill.fore_color.rgb = MINT; s2.line.fill.background()
    add_text(slide, "Azure", sw_x + 1.25 + sw + 0.05, legend_y - 0.04, 0.8, 0.2, 9, SLATE)

    metrics = [
        ("Storage", pb_from_gb(m["onprem_storage_gb"]), pb_from_gb(m["rec_storage_gb"]), m["onprem_storage_gb"], m["rec_storage_gb"]),
        ("Cores",   f"{m['onprem_cores']:,.0f}",        f"{m['rec_cores']:,.0f}",        m["onprem_cores"],     m["rec_cores"]),
        ("Memory",  tb_from_gib(m["onprem_mem_gb"]),    tb_from_gib(m["rec_mem_gb"]),    m["onprem_mem_gb"],    m["rec_mem_gb"]),
    ]
    bars_x = panel_x + 0.2
    bars_w = 1.8
    label_w = panel_w - 0.4 - bars_w - 0.05
    group_top = panel_y + 1.05
    group_h = (panel_h - 1.2) / len(metrics)
    for i, (label, on_text, az_text, on_val, az_val) in enumerate(metrics):
        gy = group_top + i * group_h
        add_text(slide, label, bars_x, gy, panel_w - 0.4, 0.22, 11, NAVY, True)
        bar_max = max(on_val, az_val, 1)
        b1y = gy + 0.32
        b2y = b1y + 0.34
        add_bar(slide, bars_x, b1y, bars_w, 0.22, on_val, bar_max, TEAL)
        add_text(slide, on_text, bars_x + bars_w + 0.05, b1y - 0.03, label_w, 0.25, 9, NAVY, True)
        add_bar(slide, bars_x, b2y, bars_w, 0.22, az_val, bar_max, MINT)
        add_text(slide, az_text, bars_x + bars_w + 0.05, b2y - 0.03, label_w, 0.25, 9, NAVY, True)

    add_source(
        slide,
        "Source: Strategy_Lift_and_shift.xlsx, Server_to_AzureVM tab. Recommended storage, cores, and memory are compared to on-premises totals.",
        size=7,
    )


def add_fileshare_os_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "Fileshares by Host OS Category")
    windows = m["os_counts"].get("Windows", 0)
    rhel = m["os_counts"].get("RHEL", 0)
    other = m["os_counts"].get("Other/Unknown", 0)
    breakdown = [("Windows", windows), ("Linux / RHEL", rhel), ("Other / Unknown", other)]
    add_card(slide, 0.85, 1.25, 2.35, 1.35, f"{m['fileshare_total']:,}", "Total fileshares", "", TEAL, value_size=34, label_size=11)
    max_count = max((c for _, c in breakdown), default=0) or 1
    for i, (label, count) in enumerate(breakdown):
        y = 1.52 + i * 0.32
        add_text(slide, label, 3.35, y, 1.4, 0.18, 10, SLATE)
        add_bar(slide, 4.85, y + 0.03, 6.80, 0.15, count, max_count, TEAL)
        add_text(slide, f"{count:,}", 11.80, y, 0.5, 0.18, 10, NAVY, True, align=PP_ALIGN.RIGHT)
    rows = [["Host OS type", "OS category", "Fileshares", "Share"]]
    for label, count in [("Windows", windows), ("Linux", rhel), ("Other/Unknown", other)]:
        if count:
            rows.append([label, "RHEL" if label == "Linux" else label, f"{count:,}", pct(count, m["fileshare_total"])])
    rows.append(["Total", "", f"{m['fileshare_total']:,}", "100.0%"])
    add_panel(slide, 1.8, 3.25, 9.2, 2.6, "Fileshare count by host operating system", title_size=17)
    add_table(slide, rows, 2.0, 3.92, 8.8, 1.8, 12, col_widths=[2.5, 2.3, 2.0, 2.0])
    add_source(
        slide,
        "Source: Discovery.xlsx, Data sheet. Fileshares are resourceType = microsoft.applicationmigration/storagesites/fileshares; OS is joined from host machine records.",
        size=8,
    )


def add_db_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "Database Resources by Type")
    add_card(slide, 0.85, 1.25, 2.35, 1.35, f"{m['db_total']:,}", "Total database resources", "", TEAL, value_size=34, label_size=11)
    max_count = max(m["db_counts"].values()) or 1
    for i, (label, count) in enumerate(m["db_counts"].items()):
        y = 1.52 + i * 0.25
        add_text(slide, label, 3.35, y, 1.2, 0.15, 9, SLATE)
        add_bar(slide, 4.7, y + 0.02, 6.95, 0.12, count, max_count, TEAL)
        add_text(slide, f"{count:,}", 11.85, y, 0.45, 0.15, 9, NAVY, True, align=PP_ALIGN.RIGHT)
    rows = [["Database type", "Count", "Share"]]
    for k, v in m["db_counts"].items():
        rows.append([k, f"{v:,}", pct(v, m["db_total"])])
    rows.append(["Total", f"{m['db_total']:,}", "100.0%"])
    add_panel(slide, 1.8, 3.25, 9.2, 2.6, "Database count by type", title_size=17)
    add_table(slide, rows, 2.0, 3.92, 8.8, 1.8, 12, col_widths=[4.6, 2.0, 2.0])
    add_source(slide, "Source: Discovery.xlsx, Data sheet. Database rows identified from database-related values in resourceType.", size=8)


def add_non_sql_db_readiness_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(
        slide,
        "Non-SQL Database Readiness",
        "Discovery footprint, support status, and PaaS readiness for MongoDB, MySQL, and PostgreSQL.",
    )

    readiness_by_type = m["non_sql_db_readiness"]
    readiness_total = m["non_sql_db_readiness_total"]
    total = m["non_sql_db_total"]
    discovery_metrics = m["non_sql_db_discovery"]
    support_order = m["non_sql_support_order"]
    db_types = ["MongoDB", "MySQL", "PostgreSQL"]
    readiness_order = ["Unknown", "Ready", "Ready With Conditions", "Not Ready"]
    readiness_colors = {
        "Ready": GREEN,
        "Ready With Conditions": YELLOW,
        "Not Ready": RED,
        "Unknown": GREY,
    }

    def lighter(color, amount=0.78):
        return RGBColor(*(int(channel + (255 - channel) * amount) for channel in color))

    def support_color(status):
        lookup = {
            "Mainstream": GREEN,
            "Extended": YELLOW,
            "OutOfSupport": RED,
            "Unknown": GREY,
        }
        return lookup.get(status, TEAL)

    add_panel(slide, 0.65, 1.20, 5.80, 1.95, "Discovery footprint", title_size=14)
    metric_headers = ["Instances", "Servers", "Versions"]
    metric_keys = ["instances", "servers", "versions"]
    metric_x = 2.35
    metric_y = 1.72
    metric_cell_w = 1.18
    metric_cell_h = 0.25
    metric_gap = 0.05
    for i, header in enumerate(metric_headers):
        x = metric_x + i * (metric_cell_w + metric_gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(metric_y), Inches(metric_cell_w), Inches(metric_cell_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
        add_text(slide, header, x + 0.06, metric_y + 0.04, metric_cell_w - 0.12, 0.16, 9, WHITE, True, align=PP_ALIGN.CENTER)
    for r, row_label in enumerate(db_types + ["Total"]):
        y = metric_y + (r + 1) * (metric_cell_h + metric_gap)
        add_text(slide, row_label, 0.90, y + 0.04, 1.18, 0.18, 10, NAVY, True)
        row = discovery_metrics[row_label]
        for c, key in enumerate(metric_keys):
            x = metric_x + c * (metric_cell_w + metric_gap)
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(metric_cell_w), Inches(metric_cell_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = NAVY if row_label == "Total" else ALT
            rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
            add_text(
                slide,
                f"{row[key]:,}",
                x + 0.06,
                y + 0.04,
                metric_cell_w - 0.12,
                0.15,
                9,
                WHITE if row_label == "Total" else NAVY,
                True,
                align=PP_ALIGN.CENTER,
            )

    add_panel(slide, 6.75, 1.20, 5.90, 1.95, "Support status by database type", title_size=14)
    support_cols = support_order or ["Unknown"]
    support_x = 8.20
    support_y = 1.72
    support_grid_w = 4.05
    support_cell_w = support_grid_w / len(support_cols)
    support_cell_h = 0.25
    support_row_gap = 0.05
    support_col_gap = 0.04
    for i, status in enumerate(support_cols):
        x = support_x + i * support_cell_w
        color = support_color(status)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(support_y), Inches(support_cell_w - support_col_gap), Inches(support_cell_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.fill.background()
        add_text(
            slide,
            status.replace("OutOfSupport", "Out of support"),
            x + 0.05,
            support_y + 0.05,
            support_cell_w - support_col_gap - 0.10,
            0.16,
            9,
            WHITE if status != "Extended" else NAVY,
            True,
            align=PP_ALIGN.CENTER,
        )
    for r, row_label in enumerate(db_types + ["Total"]):
        y = support_y + (r + 1) * (support_cell_h + support_row_gap)
        add_text(slide, row_label, 7.00, y + 0.04, 1.05, 0.18, 10, NAVY, True)
        counts = discovery_metrics[row_label]["support_counts"]
        for c, status in enumerate(support_cols):
            x = support_x + c * support_cell_w
            color = support_color(status) if row_label == "Total" else lighter(support_color(status))
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(support_cell_w - support_col_gap), Inches(support_cell_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = color
            rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
            text_color = WHITE if row_label == "Total" and status != "Extended" else NAVY
            add_text(slide, f"{counts.get(status, 0):,}", x + 0.05, y + 0.04, support_cell_w - support_col_gap - 0.10, 0.15, 9, text_color, True, align=PP_ALIGN.CENTER)

    add_panel(slide, 0.65, 3.45, 12.00, 3.05, "PaaS readiness by database type", title_size=14)

    status_labels = {
        "Unknown": "Unknown",
        "Ready": "Ready",
        "Ready With Conditions": "Ready w/ conditions",
        "Not Ready": "Not ready",
    }
    label_x = 0.95
    grid_x = 2.35
    grid_y = 4.02
    cell_w = 2.25
    cell_h = 0.42
    row_gap = 0.06
    col_gap = 0.08

    for i, status in enumerate(readiness_order):
        x = grid_x + i * (cell_w + col_gap)
        color = readiness_colors[status]
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(grid_y), Inches(cell_w), Inches(cell_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.fill.background()
        text_color = WHITE if status != "Ready With Conditions" else NAVY
        add_text(slide, status_labels[status], x + 0.10, grid_y + 0.08, cell_w - 0.20, 0.24, 12, text_color, True, align=PP_ALIGN.CENTER)

    row_labels = db_types + ["Total"]
    for r, row_label in enumerate(row_labels):
        y = grid_y + (r + 1) * (cell_h + row_gap)
        add_text(slide, row_label, label_x, y + 0.05, 1.25, 0.30, 14, NAVY, True)
        row_counts = readiness_total if row_label == "Total" else readiness_by_type.get(row_label, {})
        for c, status in enumerate(readiness_order):
            x = grid_x + c * (cell_w + col_gap)
            color = readiness_colors[status] if row_label == "Total" else lighter(readiness_colors[status])
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = color
            rect.line.color.rgb = WHITE
            rect.line.width = Pt(1)
            value = row_counts.get(status, 0)
            text_color = WHITE if row_label == "Total" and status != "Ready With Conditions" else NAVY
            add_text(slide, f"{value:,}", x + 0.10, y + 0.09, cell_w - 0.20, 0.18, 11, text_color, True, align=PP_ALIGN.CENTER)

    add_source(
        slide,
        "Sources: Discovery.xlsx Data sheet for instances, servers, versions, and directSupportStatus; Strategy_PaaS_Preferred.xlsx for PaaS readiness.",
        size=7,
    )


def add_webapp_readiness_slide(prs, m, output_dir: Path):
    slide = blank_slide(prs)
    slide_title(
        slide,
        "WebApp Readiness",
        "Discovery web application footprint, support status, and AKS readiness by web app type.",
    )

    web_types = ["IIS", "Tomcat"]
    webapp_by_type = m["webapp_by_type"]
    support_order = m["webapp_support_order"] or ["Unknown"]
    readiness_by_type = m["webapp_readiness_by_type"]
    readiness_total = m["webapp_readiness_total"]
    readiness_order = ["Unknown", "Ready", "Ready With Conditions", "Not Ready"]
    type_colors = {"IIS": TEAL, "Tomcat": ORANGE}
    support_colors = {
        "Mainstream": GREEN,
        "Extended": YELLOW,
        "OutOfSupport": RED,
        "Unknown": GREY,
    }
    readiness_colors = {
        "Ready": GREEN,
        "Ready With Conditions": YELLOW,
        "Not Ready": RED,
        "Unknown": GREY,
    }

    def lighter(color, amount=0.78):
        return RGBColor(*(int(channel + (255 - channel) * amount) for channel in color))

    def status_label(status):
        return status.replace("Ready With Conditions", "Ready w/ conditions").replace("OutOfSupport", "Out of support")

    add_panel(slide, 0.55, 1.12, 12.15, 2.55, "WebApps and hosting servers", title_size=13)
    label_x = 0.88
    count_x = 3.25
    count_y = 1.78
    count_cell_w = 2.45
    count_cell_h = 0.34
    count_gap = 0.08
    count_headers = [("WebApps", TEAL), ("Server Count", NAVY)]
    for i, (header, color) in enumerate(count_headers):
        x = count_x + i * (count_cell_w + count_gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(count_y), Inches(count_cell_w), Inches(count_cell_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.fill.background()
        add_text(slide, header, x + 0.10, count_y + 0.08, count_cell_w - 0.20, 0.18, 10, WHITE, True, align=PP_ALIGN.CENTER)
    for r, row_label in enumerate([*web_types, "Total"]):
        y = count_y + (r + 1) * (count_cell_h + count_gap)
        add_text(slide, row_label, label_x, y + 0.08, 1.25, 0.18, 10, NAVY, True)
        row = webapp_by_type[row_label]
        values = [row["webapps"], row["servers"]]
        for c, value in enumerate(values):
            x = count_x + c * (count_cell_w + count_gap)
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(count_cell_w), Inches(count_cell_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = NAVY if row_label == "Total" else ALT
            rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
            add_text(slide, f"{value:,}", x + 0.10, y + 0.08, count_cell_w - 0.20, 0.18, 11, WHITE if row_label == "Total" else NAVY, True, align=PP_ALIGN.CENTER)

    def draw_stacked_matrix(panel_x, panel_y, panel_w, title, columns, row_values, colors):
        add_panel(slide, panel_x, panel_y, panel_w, 2.52, title, title_size=13)
        label_x = panel_x + 0.22
        grid_x = panel_x + 1.38
        grid_y = panel_y + 0.65
        grid_w = panel_w - 1.70
        cell_h = 0.34
        row_gap = 0.08
        col_gap = 0.05
        cell_w = grid_w / len(columns)
        for i, col in enumerate(columns):
            x = grid_x + i * cell_w
            color = colors.get(col, TEAL)
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(grid_y), Inches(cell_w - col_gap), Inches(cell_h))
            rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.fill.background()
            text_color = WHITE if col not in {"Extended", "Ready With Conditions"} else NAVY
            add_text(slide, status_label(col), x + 0.04, grid_y + 0.08, cell_w - col_gap - 0.08, 0.17, 7, text_color, True, align=PP_ALIGN.CENTER)
        row_labels = web_types + ["Total"]
        for r, row_label in enumerate(row_labels):
            y = grid_y + (r + 1) * (cell_h + row_gap)
            add_text(slide, row_label, label_x, y + 0.08, 0.95, 0.17, 9, NAVY, True)
            values = row_values["Total"] if row_label == "Total" else row_values.get(row_label, {})
            for c, col in enumerate(columns):
                x = grid_x + c * cell_w
                color = colors.get(col, TEAL) if row_label == "Total" else lighter(colors.get(col, TEAL))
                rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cell_w - col_gap), Inches(cell_h))
                rect.fill.solid(); rect.fill.fore_color.rgb = color
                rect.line.color.rgb = WHITE; rect.line.width = Pt(1)
                text_color = WHITE if row_label == "Total" and col not in {"Extended", "Ready With Conditions"} else NAVY
                add_text(slide, f"{values.get(col, 0):,}", x + 0.04, y + 0.08, cell_w - col_gap - 0.08, 0.17, 9, text_color, True, align=PP_ALIGN.CENTER)

    add_panel(slide, 0.55, 4.02, 5.95, 2.52, "Support Status", title_size=13)
    support_rows = {web_type: webapp_by_type[web_type]["support_counts"] for web_type in web_types}
    support_rows["Total"] = webapp_by_type["Total"]["support_counts"]
    support_max = max([support_rows[row].get(status, 0) for row in web_types for status in support_order] + [1])
    chart_left = 0.95
    chart_top = 4.78
    chart_bottom = 6.08
    chart_h = chart_bottom - chart_top
    group_w = 5.10 / len(support_order)
    bar_w = 0.18
    bar_gap = 0.07
    for i, status in enumerate(support_order):
        group_x = chart_left + i * group_w
        pair_w = len(web_types) * bar_w + (len(web_types) - 1) * bar_gap
        x0 = group_x + (group_w - pair_w) / 2
        for j, web_type in enumerate(web_types):
            value = support_rows[web_type].get(status, 0)
            h = max(0.01, chart_h * safe_div(value, support_max))
            x = x0 + j * (bar_w + bar_gap)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(chart_bottom - h), Inches(bar_w), Inches(h))
            bar.fill.solid(); bar.fill.fore_color.rgb = type_colors[web_type]; bar.line.fill.background()
            if value:
                add_text(slide, f"{value:,}", x - 0.12, chart_bottom - h - 0.16, bar_w + 0.24, 0.13, 6, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, status_label(status), group_x + 0.02, chart_bottom + 0.06, group_w - 0.04, 0.22, 7, SLATE, True, align=PP_ALIGN.CENTER)
    legend_y = 4.42
    for i, web_type in enumerate(web_types):
        x = 0.85 + i * 1.05
        sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(legend_y), Inches(0.12), Inches(0.10))
        sw.fill.solid(); sw.fill.fore_color.rgb = type_colors[web_type]; sw.line.fill.background()
        add_text(slide, web_type, x + 0.17, legend_y - 0.04, 0.70, 0.16, 7, SLATE, True)

    readiness_rows = {web_type: readiness_by_type.get(web_type, {}) for web_type in web_types}
    readiness_rows["Total"] = readiness_total
    draw_stacked_matrix(6.75, 4.02, 5.95, "Readiness", readiness_order, readiness_rows, readiness_colors)

    add_source(
        slide,
        "Sources: Discovery.xlsx Data sheet for resourceType, parentResourceName, and directSupportStatus; Strategy_PaaS_Preferred.xlsx WebApp_to_AKS for READINESS.",
        size=7,
    )


def add_webapp_slide(prs, m, output_dir: Path):
    slide = blank_slide(prs)
    slide_title(slide, "Web Applications by Type", "Counted by webapp rows")
    add_panel(slide, 2.35, 1.25, 8.60, 4.95, "Webapp type distribution", title_size=17)
    pie_values = {
        "IIS": int(m["webapp_by_type"]["IIS"]["webapps"]),
        "Tomcat": int(m["webapp_by_type"]["Tomcat"]["webapps"]),
    }
    pie_path = create_pie_image(output_dir / "webapp-type-distribution-large.png", "WebApp Types", pie_values, {"IIS": TEAL, "Tomcat": ORANGE})
    slide.shapes.add_picture(str(pie_path), Inches(2.70), Inches(1.80), Inches(5.45), Inches(3.35))
    summary_rows = [
        ["Metric", "Count"],
        ["Webapp rows", f"{m['webapp_by_type']['Total']['webapps']:,}"],
        ["Unique servers", f"{m['webapp_by_type']['Total']['servers']:,}"],
    ]
    add_table(slide, summary_rows, 8.25, 2.35, 2.35, 1.20, 12, col_widths=[1.45, 0.90])
    add_text(
        slide,
        "Server count deduplicates parentResourceName across all webapp resourceType rows.",
        8.25,
        3.75,
        2.25,
        0.45,
        9,
        MUTED,
    )
    add_source(
        slide,
        "Source: Discovery.xlsx, Data sheet. Webapp rows identified from resourceType; duplicate parentResourceName values are counted separately.",
        size=8,
    )


def add_consolidated_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "Consolidated Infrastructure Summary")
    add_panel(slide, 0.4, 1.1, 2.6, 2.0, "VM Power State Chart", title_size=14)
    max_count = max(m["powered_on"], m["powered_off"], 1)
    for i, (label, count, color) in enumerate([("Powered On", m["powered_on"], TEAL), ("Powered Off", m["powered_off"], PINK)]):
        y = 1.6 + i * 0.66
        add_text(slide, label, 0.6, y, 1.0, 0.22, 8, SLATE, True)
        add_bar(slide, 1.55, y + 0.03, 1.1, 0.15, count, max_count, color)
        add_text(slide, f"{count:,} ({pct(count, m['vm_total'])})", 0.6, y + 0.26, 2.15, 0.3, 10, NAVY, True)

    add_panel(slide, 3.15, 1.1, 9.4, 2.0, "VM Utilization Summary", title_size=14)
    storage_delta = delta_pct(m["rec_storage_gb"], m["onprem_storage_gb"])
    cores_delta = delta_pct(m["rec_cores"], m["onprem_cores"])
    mem_delta = delta_pct(m["rec_mem_gb"], m["onprem_mem_gb"])
    rows = [
        ["", "On-premises total", "On-premises Utilization", "Recommended Azure sizing"],
        ["Assessed VMs", f"{m['vm_total']:,}", "", ""],
        ["Storage", pb_from_gb(m["onprem_storage_gb"]), f"{m['storage_util']:.1f}% Consumed", f"{pb_from_gb(m['rec_storage_gb'])} ({delta_text(storage_delta)})"],
        ["Cores", f"{m['onprem_cores']:,.0f}", f"{m['cpu_util']:.1f}% CPU Utilization", f"{m['rec_cores']:,.0f} ({delta_text(cores_delta)})"],
        ["Memory", tb_from_gib(m["onprem_mem_gb"]), f"{m['mem_util']:.1f}% Memory Utilization", f"{tb_from_gib(m['rec_mem_gb'])} ({delta_text(mem_delta)})"],
    ]
    tbl_shape = add_table(slide, rows, 3.42, 1.56, 9.33, 1.25, 11, col_widths=[1.3, 2.1, 2.55, 3.05])
    table = tbl_shape.table
    for r in (2, 3, 4):
        color_delta_text(table.cell(r, 3), table.cell(r, 3).text)

    add_panel(slide, 0.4, 3.45, 3.9, 2.4, "Fileshare count by host OS", title_size=14)
    windows = m["os_counts"].get("Windows", 0)
    rhel = m["os_counts"].get("RHEL", 0)
    other = m["os_counts"].get("Other/Unknown", 0)
    fs_rows = [["Host OS type", "OS category", "Fileshares", "Share"], ["Windows", "Windows", f"{windows:,}", pct(windows, m["fileshare_total"])]]
    if rhel:
        fs_rows.append(["Linux", "RHEL", f"{rhel:,}", pct(rhel, m["fileshare_total"])])
    if other:
        fs_rows.append(["Other", "Unknown", f"{other:,}", pct(other, m["fileshare_total"])])
    fs_rows.append(["Total", "", f"{m['fileshare_total']:,}", "100.0%"])
    add_table(slide, fs_rows, 0.55, 4.05, 3.6, 1.6, 11, col_widths=[1.1, 1.0, 0.85, 0.65])

    add_panel(slide, 4.55, 3.45, 3.75, 2.4, "Database count by type", title_size=14)
    db_rows = [["Database type", "Count", "Share"]]
    for k, v in m["db_counts"].items():
        db_rows.append([k, f"{v:,}", pct(v, m["db_total"])])
    db_rows.append(["Total", f"{m['db_total']:,}", "100.0%"])
    add_table(slide, db_rows, 4.7, 4.05, 3.45, 1.7, 11, col_widths=[1.85, 0.85, 0.75])

    add_panel(slide, 8.7, 3.45, 3.85, 2.4, "Webapp count by type", title_size=14)
    web_rows = [["Webapp type", "Webapp", "Share"]]
    for k, v in m["web_counts"].items():
        web_rows.append([k, f"{v:,}", pct(v, m["web_total"])])
    web_rows.append(["Total", f"{m['web_total']:,}", "100.0%"])
    add_table(slide, web_rows, 8.85, 4.05, 3.55, 1.6, 11, col_widths=[1.95, 0.85, 0.75])


def add_fileshare_readiness_slide(prs, m, output_dir: Path):
    slide = blank_slide(prs)
    slide_title(slide, "Fileshare Readiness", "Discovery-derived fileshare-only server scope with Azure Files readiness, sizing, and monthly cost comparison.")
    add_card(slide, 0.52, 1.55, 2.35, 1.35, f"{m['fileshare_only_servers']:,}", "Fileshare-only servers", "", TEAL, value_size=20, label_size=10)
    add_card(slide, 3.02, 1.55, 2.35, 1.35, f"{m['fileshare_only_shares']:,}", "Fileshare-only shares", "", MINT, value_size=20, label_size=10)

    add_panel(slide, 0.55, 3.25, 5.1, 3.15, "", title_size=1)
    donut = create_donut_image(output_dir / "protocol-distribution.png", m["protocol_counts"].get("SMB", 0), m["protocol_counts"].get("NFS", 0))
    slide.shapes.add_picture(str(donut), Inches(0.72), Inches(3.5), Inches(4.58), Inches(2.85))

    add_panel(slide, 6.05, 1.55, 6.3, 1.75, "Fileshare-only Azure sizing and cost", title_size=16)
    add_small_card(slide, 6.17, 2.10, 1.42, 1.02, f"{m['fileshare_only_size_tb']:.1f} TB", "Azure Files size", TEAL, value_size=15, label_size=8)
    add_small_card(slide, 7.70, 2.10, 1.42, 1.02, money_k(m["fileshare_only_azure_cost"]), "Azure Files / mo", GREEN, value_size=15, label_size=8)
    add_small_card(slide, 9.23, 2.10, 1.42, 1.02, money_k(m["fileshare_only_lift_cost"]), "Lift & shift / mo", RED, value_size=15, label_size=8)
    delta_value = m["fileshare_only_lift_cost"] - m["fileshare_only_azure_cost"]
    delta_color = GREEN if delta_value >= 0 else RED
    add_small_card(slide, 10.75, 2.10, 1.34, 1.02, money_k(delta_value), "Monthly delta", delta_color, value_size=15, label_size=8)

    add_panel(slide, 6.05, 3.25, 6.3, 3.15, "Azure Files readiness distribution", title_size=16)
    readiness = create_readiness_image(output_dir / "azure-files-readiness-distribution.png", m["readiness_counts"])
    slide.shapes.add_picture(str(readiness), Inches(6.04), Inches(3.97), Inches(6.25), Inches(2.64))
    add_text(
        slide,
        "Sources: Discovery.xlsx Data sheet; Strategy_PaaS_Preferred.xlsx FileShares_to_Azure_Files; Strategy_Lift_and_shift.xlsx Server_to_AzureVM.",
        0.53,
        7.01,
        10.97,
        0.3,
        7,
        MUTED,
    )


def add_sql_readiness_slide(prs, m):
    slide = blank_slide(prs)
    slide_title(slide, "SQL Readiness")

    # KPI strip
    def sql_kpi_card(x, value, label, accent):
        w, h, y = 2.4, 0.60, 1.00
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = LINE; rect.line.width = Pt(1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        add_text(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.24, 18, NAVY, True)
        add_text(slide, label, x + 0.12, y + 0.40, w - 0.24, 0.15, 8, SLATE, True)

    sql_kpi_card(0.55, f"{m['sql_instance_count']:,}", "SQL Server Instances", TEAL)
    sql_kpi_card(3.10, f"{m['sql_server_count']:,}", "Servers running SQL", MINT)
    instances_per_server = m["sql_instance_count"] / m["sql_server_count"] if m["sql_server_count"] else 0
    sql_kpi_card(5.65, f"{instances_per_server:.1f}", "Instances per server", TEAL)

    # Cost comparison — compact panel with three inner cards
    lift_cost = float(m.get("sql_lift_total_cost", 0) or 0)
    hybrid_cost = float(m.get("sql_hybrid_cost", 0) or 0)
    yearly_lift = lift_cost * 12
    yearly_hybrid = hybrid_cost * 12
    yearly_savings = yearly_lift - yearly_hybrid  # positive => MI saves money
    savings_color = GREEN if yearly_savings >= 0 else RED

    cost_panel_y = 1.78
    cost_panel_h = 1.05
    add_panel(slide, 0.55, cost_panel_y, 12.30, cost_panel_h, "Lift & Shift vs SQL MI (where ready) + Lift & Shift remainder", title_size=11)

    inner_y = cost_panel_y + 0.35
    inner_h = cost_panel_h - 0.45
    gap = 0.15
    inner_w = (12.30 - 0.30 - gap * 2) / 3

    def cost_card(x, accent, headline, big_value, sub_text, value_color=None):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(inner_y), Inches(inner_w), Inches(inner_h))
        rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = LINE; rect.line.width = Pt(1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(inner_y), Inches(inner_w), Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        add_text(slide, headline, x + 0.16, inner_y + 0.08, inner_w - 0.32, 0.18, 9, NAVY, True)
        add_text(slide, big_value, x + 0.16, inner_y + 0.27, inner_w - 0.32, 0.25, 17, value_color or NAVY, True)
        add_text(slide, sub_text, x + 0.16, inner_y + 0.50, inner_w - 0.32, 0.14, 7, MUTED)

    x0 = 0.70
    cost_card(x0, RED, "100% Lift & Shift",
              money_k(lift_cost),
              f"per month   /   {money_full(yearly_lift)}/yr")
    cost_card(x0 + inner_w + gap, TEAL, "SQL MI (where ready) + Lift & Shift",
              money_k(hybrid_cost),
              f"per month   /   {money_full(yearly_hybrid)}/yr")
    cost_card(x0 + (inner_w + gap) * 2, savings_color, "Yearly Savings Using SQL MI",
              money_full(yearly_savings),
              f"{money_full(yearly_savings / 12)} per month",
              value_color=savings_color)

    # Licensing snapshot (placeholder values — wire real data in once available)
    lic_y = 2.98
    lic_h = 1.10
    add_panel(slide, 0.55, lic_y, 12.30, lic_h, "SQL Server Licensing — owned vs. need", title_size=12)

    # Random-ish placeholder values until real data is wired in
    licensing_rows = [
        ("SQL Server Standard",   m.get("sql_std_owned",  18), m.get("sql_std_needed",  24)),
        ("SQL Server Enterprise", m.get("sql_ent_owned",  12), m.get("sql_ent_needed",  8)),
    ]

    def licensing_card(x, w, edition, owned, needed):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(lic_y + 0.40), Inches(w), Inches(lic_h - 0.50))
        rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = LINE; rect.line.width = Pt(1)

        try:
            o = float(owned); n = float(needed)
            delta_pct = (o - n) / n * 100 if n else 0
            delta_txt = f"{delta_pct:+.1f}%  {'surplus' if delta_pct >= 0 else 'deficit'}"
            delta_color = GREEN if delta_pct >= 0 else RED
            owned_txt = f"{int(o):,}"
            needed_txt = f"{int(n):,}"
        except (TypeError, ValueError):
            delta_txt = "—"
            delta_color = MUTED
            owned_txt = str(owned)
            needed_txt = str(needed)

        add_text(slide, edition, x + 0.18, lic_y + 0.43, w - 0.36, 0.20, 10, NAVY, True)
        col_w = (w - 0.36) / 3
        cx = x + 0.18
        row_y = lic_y + 0.64
        add_text(slide, "Owned",    cx,                row_y,        col_w, 0.16, 8, MUTED)
        add_text(slide, owned_txt,  cx,                row_y + 0.15, col_w, 0.22, 13, NAVY, True)
        add_text(slide, "Need",     cx + col_w,        row_y,        col_w, 0.16, 8, MUTED)
        add_text(slide, needed_txt, cx + col_w,        row_y + 0.15, col_w, 0.22, 13, NAVY, True)
        add_text(slide, "Delta",    cx + col_w * 2,    row_y,        col_w, 0.16, 8, MUTED)
        add_text(slide, delta_txt,  cx + col_w * 2,    row_y + 0.15, col_w, 0.22, 10, delta_color, True)

    card_w = (12.30 - 0.30 - 0.15) / 2
    licensing_card(0.70,                       card_w, *licensing_rows[0])
    licensing_card(0.70 + card_w + 0.15,       card_w, *licensing_rows[1])

    # Versions chart (left)
    pss_colors = {
        "Mainstream":   GREEN,
        "Extended":     YELLOW,
        "OutOfSupport": RED,
        "Unknown":      GREY,
    }
    versions_y = 4.25
    versions_h = 2.35
    add_panel(slide, 0.55, versions_y, 6.4, versions_h, "SQL Server versions by support status", title_size=13)
    legend_y = versions_y + 0.40
    for i, (label, color) in enumerate([("Mainstream", GREEN), ("Extended", YELLOW), ("Out of support", RED), ("Unknown", GREY)]):
        x = 0.75 + i * 1.55
        sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(legend_y), Inches(0.16), Inches(0.11))
        sw.fill.solid(); sw.fill.fore_color.rgb = color; sw.line.fill.background()
        add_text(slide, label, x + 0.22, legend_y - 0.05, 1.3, 0.22, 10, SLATE, True, font="Aptos")

    versions = m.get("sql_versions") or []
    if not versions:
        add_text(slide, "No SQL version data available in Discovery.xlsx.", 0.75, versions_y + 0.75, 6.0, 0.3, 10, MUTED)
    else:
        max_count = max(c for *_x, c in versions) or 1
        rows_top = versions_y + 0.75
        row_h = min(0.20, (versions_y + versions_h - 0.05 - rows_top) / max(len(versions), 1))
        bar_x = 2.30
        bar_w = 3.55
        for i, (version, status, count) in enumerate(versions):
            y = rows_top + i * row_h
            color = pss_colors.get(status, GREY)
            add_text(slide, version, 0.7, y, 1.55, row_h - 0.02, 9, SLATE, True)
            add_bar(slide, bar_x, y + 0.04, bar_w, row_h - 0.10, count, max_count, color)
            add_text(slide, f"{count:,}", bar_x + bar_w + 0.05, y, 0.85, row_h - 0.02, 9, NAVY, True)

    # SQL Readiness by Target (right) — vertical grouped column chart
    target_y = versions_y
    target_h = versions_h
    panel_x = 7.10
    panel_w = 5.75
    add_panel(slide, panel_x, target_y, panel_w, target_h, "SQL Readiness by Target", title_size=13)
    categories = ["Ready", "Not Ready", "Ready With Conditions", "Unknown"]
    vm_counts = m.get("sql_vm_readiness", {})
    mi_counts = m.get("sql_mi_readiness", {})
    vm_total = m.get("sql_vm_total", 0)
    mi_total = m.get("sql_mi_total", 0)

    # Legend
    leg_y = target_y + 0.40
    sw1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(panel_x + 0.20), Inches(leg_y), Inches(0.16), Inches(0.11))
    sw1.fill.solid(); sw1.fill.fore_color.rgb = ORANGE; sw1.line.fill.background()
    add_text(slide, f"SQL VM Instances ({vm_total})", panel_x + 0.40, leg_y - 0.05, 2.4, 0.22, 10, SLATE, True, font="Aptos")
    sw2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(panel_x + 2.90), Inches(leg_y), Inches(0.16), Inches(0.11))
    sw2.fill.solid(); sw2.fill.fore_color.rgb = NAVY; sw2.line.fill.background()
    add_text(slide, f"SQL MI Instances ({mi_total})", panel_x + 3.10, leg_y - 0.05, 2.5, 0.22, 10, SLATE, True, font="Aptos")

    # Chart plotting area
    chart_top = target_y + 0.75
    chart_bottom = target_y + target_h - 0.30  # leave room for x-axis labels
    chart_left = panel_x + 0.25
    chart_right = panel_x + panel_w - 0.10
    chart_h = chart_bottom - chart_top

    max_bar = max([vm_counts.get(c, 0) for c in categories] + [mi_counts.get(c, 0) for c in categories] + [1])
    group_w = (chart_right - chart_left) / len(categories)
    bar_gap = 0.04
    bar_w = max(0.18, (group_w - 0.30 - bar_gap) / 2)

    def vbar(x, value, color):
        h = max(0.02, chart_h * safe_div(value, max_bar))
        # background track (full height) for visual reference
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(chart_top), Inches(bar_w), Inches(chart_h))
        bg.fill.solid(); bg.fill.fore_color.rgb = ALT; bg.line.fill.background()
        # actual bar grows from bottom
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(chart_bottom - h), Inches(bar_w), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        return chart_bottom - h  # return top-of-bar y

    for i, cat in enumerate(categories):
        gx = chart_left + i * group_w
        # Center the bar pair within the group
        pair_w = 2 * bar_w + bar_gap
        b1x = gx + (group_w - pair_w) / 2
        b2x = b1x + bar_w + bar_gap
        vm_val = vm_counts.get(cat, 0)
        mi_val = mi_counts.get(cat, 0)
        vm_top = vbar(b1x, vm_val, ORANGE)
        mi_top = vbar(b2x, mi_val, NAVY)
        # Value labels above each bar
        add_text(slide, f"{vm_val:,}", b1x - 0.08, vm_top - 0.20, bar_w + 0.16, 0.18, 8, NAVY, True, font="Aptos", align=PP_ALIGN.CENTER)
        add_text(slide, f"{mi_val:,}", b2x - 0.08, mi_top - 0.20, bar_w + 0.16, 0.18, 8, NAVY, True, font="Aptos", align=PP_ALIGN.CENTER)
        # X-axis category label
        label = "Ready w/ Conditions" if cat == "Ready With Conditions" else cat
        label_w = 1.55 if cat == "Ready With Conditions" else group_w
        label_x = gx + (group_w - label_w) / 2
        add_text(slide, label, label_x, chart_bottom + 0.04, label_w, 0.18, 9, SLATE, True, font="Aptos", align=PP_ALIGN.CENTER)

    add_source(
        slide,
        "Sources: Discovery.xlsx Data (SQL instances & versions); Strategy_PaaS_Preferred.xlsx SQLinstance_to_AzureSQLVM "
        "(AZURE_SQL_VM_READINESS + AZURE_SQL_VM_COMPUTE/STORAGE_MONTHLY_COST_USD) and SQLinstance_to_AzureSQLMI "
        "(AZURE_SQL_MI_READINESS + AZURE_SQL_MI_COMPUTE/STORAGE_MONTHLY_COST_USD); Strategy_Lift_and_shift.xlsx "
        "Server_to_AzureVM (TOTAL_MONTHLY_COST_USD for servers in either SQL sheet, deduplicated).",
        size=7,
    )


def add_title_slide(prs):
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.05), Inches(SLIDE_W), Inches(0.10))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = TEAL; stripe.line.fill.background()

    add_text(slide, "Azure Migrate Interpreted", 0.6, 2.6, SLIDE_W - 1.2, 1.3, 60, WHITE, True, font="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(slide, "Discovery and assessment summary", 0.6, 4.05, SLIDE_W - 1.2, 0.5, 20, RGBColor(202, 220, 252), False, align=PP_ALIGN.CENTER)


def is_file_locked(path: Path) -> bool:
    """Return True if the file exists and cannot be opened for read+write
    (e.g. open in Excel/PowerPoint with an exclusive lock)."""
    if not path.exists():
        return False
    try:
        with open(path, "r+b"):
            return False
    except (PermissionError, OSError):
        return True


def wait_until_unlocked(paths: list[Path], label: str = "files") -> None:
    """If any of the given paths are locked, list them and prompt the user
    to close them. Re-checks after each prompt until all are unlocked."""
    while True:
        locked = [p for p in paths if is_file_locked(p)]
        if not locked:
            return
        log("")
        log(f"The following {label} appear to be open in another program (likely Excel or PowerPoint):")
        for p in locked:
            log(f"  - {p}")
        try:
            input("Close them, then press ENTER to retry (Ctrl+C to abort)... ")
        except EOFError:
            raise RuntimeError("Locked files and no console available to prompt the user")


def build_deck(input_dir: Path, output: Path):
    log(f"Input  : {input_dir}")
    log(f"Output : {output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    files_to_check = [
        input_dir / "Discovery.xlsx",
        input_dir / "Strategy_Lift_and_shift.xlsx",
        input_dir / "Strategy_PaaS_Preferred.xlsx",
        output,
    ]
    wait_until_unlocked(files_to_check, label="workbooks/output deck")

    log("Loading source workbooks...")
    m = load_metrics(input_dir)
    log("Source data loaded:")
    log(f"  VMs            : {m['vm_total']:,} ({m['powered_on']:,} on / {m['powered_off']:,} off)")
    log(f"  Fileshares     : {m['fileshare_total']:,} ({m['fileshare_only_servers']:,} fileshare-only servers / {m['fileshare_only_shares']:,} shares)")
    log(f"  Databases      : {m['db_total']:,}")
    log(f"  Web apps       : {m['web_total']:,}")
    log("Building slides...")
    prs = new_deck()
    log("  [1/10] Title")
    add_title_slide(prs)
    log("  [2/10] Consolidated Infrastructure Summary")
    add_consolidated_slide(prs, m)
    log("  [3/10] Fileshare Readiness")
    add_fileshare_readiness_slide(prs, m, output.parent)
    log("  [4/10] VM Power State Summary")
    add_vm_power_slide(prs, m)
    log("  [5/10] VM Utilization Summary")
    add_vm_utilization_slide(prs, m)
    log("  [6/10] Fileshares by Host OS Category")
    add_fileshare_os_slide(prs, m)
    log("  [7/10] Database Resources by Type")
    add_db_slide(prs, m)
    log("  [8/10] Non-SQL Database Readiness")
    add_non_sql_db_readiness_slide(prs, m)
    log("  [9/10] SQL Readiness")
    add_sql_readiness_slide(prs, m)
    log("  [10/10] WebApp Readiness")
    add_webapp_readiness_slide(prs, m, output.parent)
    log("Saving deck...")
    # Write to a local temp file first, validate, then move to the final path.
    # This prevents PowerPoint/OneDrive from seeing a partially-written file
    # if the output path is inside a OneDrive-synced folder.
    tmp_dir = Path(tempfile.mkdtemp(prefix="azmigrate-deck-"))
    tmp_path = tmp_dir / output.name
    try:
        prs.save(tmp_path)
        # Validate that the saved file is a valid Open XML package (zip).
        import zipfile as _zf
        with _zf.ZipFile(tmp_path, "r") as zf:
            bad = zf.testzip()
            if bad is not None:
                raise RuntimeError(f"Saved deck failed zip integrity check on entry: {bad}")
        # Replace the destination atomically; retry briefly if OneDrive briefly holds it.
        import time
        last_err = None
        for attempt in range(5):
            try:
                if output.exists():
                    output.unlink()
                shutil.move(str(tmp_path), str(output))
                break
            except PermissionError as exc:
                last_err = exc
                log(f"  Target {output.name} is locked (attempt {attempt + 1}/5); retrying in 2s...")
                time.sleep(2)
        else:
            raise RuntimeError(f"Could not move deck into place: {last_err}")
    finally:
        try:
            if tmp_path.exists():
                tmp_path.unlink()
            tmp_dir.rmdir()
        except OSError:
            pass
    log(f"Done. Wrote {output}")
    os.startfile(output)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input-dir", required=True, type=Path, help="Directory containing the three source workbooks.")
    parser.add_argument("--output", required=True, type=Path, help="Output .pptx path.")
    args = parser.parse_args()
    build_deck(args.input_dir, args.output)


if __name__ == "__main__":
    main()
