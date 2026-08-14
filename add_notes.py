"""
Inject business-assessment notes into each slide of the Kettering deck.
Reads _kettering_repair_tmp.pptx (already de-MIP'd) and writes
Kettering-Slides-with-notes.pptx with Notes populated per slide.
"""
from pptx import Presentation

SRC = "_kettering_repair_tmp.pptx"
OUT = "Kettering-Slides-with-notes.pptx"

NOTES = {
    1: None,  # Title slide - skip
    2: """CONSOLIDATED INFRASTRUCTURE SUMMARY - ASSESSMENT

RIGHT-SIZING OPPORTUNITY: Recommended Azure sizing shows storage down 48.5% (3.03 PB to 1.56 PB), cores down 23.4% (8,821 to 6,758), and memory down 16.0% (44.4 TB to 37.3 TB). The on-prem estate is significantly overprovisioned relative to actual utilization (48.8% storage, 48.9% CPU, 32.9% memory) - the customer pays for capacity it isn't using today. This benefit applies broadly, independent of any PaaS conversion.

VM RETIREMENT IMPACT: Across the workload-specific slides in this deck: 475 dedicated WebApp servers, 99 SQL-only servers, and 37 non-SQL database servers are candidates for full decommission via PaaS conversion (App Service/AKS, SQL MI, and Cosmos DB/Flexible Server respectively). That is 611+ VMs - over 40% of the 1,472-VM estate - eligible for full retirement, not just right-sizing. Each retired VM removes its OS/software licensing, patching cycles, backup jobs, monitoring configuration, and ops labor.

QUICK WIN: 74 powered-off VMs (5.0%) are consuming storage, backup, and licensing overhead while idle. Confirm with workload owners and decommission before the broader migration wave begins - near-zero-effort cost recovery.

SECURITY RISK ROLL-UP: 111 WebApp servers (27 confirmed out-of-support + 84 unknown), 14 of 37 non-SQL DB instances (38%) out of support, and 3 SQL 2014 instances plus 186 more (74% of SQL fleet) on an expiring support clock. This is a portfolio-wide risk pattern, not isolated to one workload.

OTHER OBSERVATIONS: Database estate is SQL Server-dominant (87.2%) - the single largest modernization lever. Webapp estate is IIS-dominant (99.5%) - enables a highly standardized migration playbook. Fileshares are 99.6% Windows - strong Azure Files candidate, potentially eliminating dedicated file servers. 1,472 VMs support only 746 fileshares + 289 databases + 870 webapps combined, meaning many VMs are mixed-role - a full dependency-mapping pass would likely expand the retirement count further.

AI OPPORTUNITY: Microsoft Fabric can unify SQL Server, MongoDB, MySQL, PostgreSQL, and fileshare data into a single governed estate once migrated - the prerequisite for enterprise Copilot/AI adoption. Copilot in Azure SQL, Copilot in Fabric, and M365 Copilot become viable once data exits fragmented, partially out-of-support on-prem silos. Azure Advisor/Copilot for Azure continues right-sizing post-migration on an ongoing basis. Azure AI Search + Azure OpenAI can ground RAG patterns directly against this consolidated data for new customer/employee-facing AI experiences. This reframes the project from pure infrastructure migration into a data-and-AI enablement initiative worth executive-level visibility.""",

    3: """FILESHARE READINESS - ASSESSMENT

VM RETIREMENT IMPACT: 68 fileshare-only servers (hosting 275 shares, 147.2 TB) are candidates for full decommission if migrated to Azure Files. These servers exist purely to serve files - once moved, the customer eliminates OS licensing, patch cycles, backup infrastructure, and file-server admin overhead for all 68 boxes, not just a cost/size reduction.

COST IMPACT: Azure Files sizing shows $11.3K/mo vs a projected $21.2K/mo for lift-and-shift VMs - a $9.9K/mo delta ($118.8K/yr) in favor of PaaS, on top of the VM retirement savings above. This is a direct, quantifiable number for the business case.

SECURITY/OPERATIONAL RISK: File servers are a common ransomware target (mass file encryption). Retiring 68 dedicated file servers in favor of Azure Files (with built-in redundancy, soft-delete, and snapshot capabilities) meaningfully reduces this attack surface versus maintaining on-prem file server patching cadence.

OTHER OBSERVATIONS: This is one of the cleanest ROI stories in the deck - concrete monthly cost delta plus full VM elimination for a defined, isolated server population (fileshare-only, no mixed-role complexity). Good candidate for an early migration wave to build momentum/confidence before tackling higher-complexity workloads (SQL, WebApp).

AI OPPORTUNITY: Once file data lives in Azure Files, it becomes eligible for Azure AI Search indexing and Microsoft 365 Copilot grounding (if migrated further into SharePoint/OneDrive), enabling natural-language search and summarization across previously siloed on-prem file shares - content that today is invisible to any AI tooling.""",

    4: """VM POWER STATE SUMMARY - ASSESSMENT

QUICK WIN - IMMEDIATE COST RECOVERY: 74 VMs (5.0%) are powered off today, meaning they are consuming storage, backup, and licensing/inventory overhead while providing zero business value. This should be the very first action item in the engagement - confirm with workload owners that these are truly abandoned, then decommission immediately. This delivers savings before any migration wave begins and requires no re-platforming effort.

SECURITY RISK: Powered-off VMs are frequently the least-patched, least-monitored assets in an environment - if any are ever reactivated without a patching review, they represent a silent security gap. Treat unresolved power-off status as a data-hygiene risk, not just a cost issue.

OTHER OBSERVATIONS: 1,398 powered-on VMs (95.0%) represent the active migration scope for planning purposes - the 74 powered-off machines should be excluded from migration wave sizing/cost estimates once confirmed for decommission, so they don't inflate the perceived migration scope or cost.

AI OPPORTUNITY: Azure Advisor and Copilot for Azure can continuously flag underutilized or idle VMs post-migration (not just at a point-in-time assessment like this one), turning this one-time "74 idle VMs" finding into an ongoing automated governance capability rather than a manual audit exercise.""",

    5: """VM UTILIZATION SUMMARY - ASSESSMENT

RIGHT-SIZING IMPACT: Recommended Azure sizing reduces storage by 48.5% (3.03 PB to 1.56 PB), cores by 23.4% (8,821 to 6,758), and memory by 16.0% (44.4 TB to 37.3 TB) versus current on-premises footprint. Current utilization (48.8% storage, 48.9% CPU, 32.9% memory) confirms significant overprovisioning today - the customer is paying for capacity that sits idle, and this gap converts directly into cost savings once right-sized in Azure.

VM RETIREMENT IMPACT: This slide reflects the full 1,472-VM assessed footprint before workload-specific PaaS conversion. Combined with the SQL, WebApp, and Non-SQL DB readiness findings elsewhere in this deck, a meaningful share of this footprint (600+ VMs) can move beyond right-sizing into full retirement via PaaS - compounding the storage/core/memory reduction shown here with elimination of the licensing and operational overhead tied to those servers entirely.

OTHER OBSERVATIONS: A 48.5% storage reduction is unusually large and worth validating with the customer - it may reflect thin-provisioning cleanup, deduplication/compression benefits in Azure, or genuine over-allocation on-prem. Any of these are a legitimate win, but the underlying driver should be understood before quoting it to stakeholders.

AI OPPORTUNITY: Post-migration, Azure Monitor + Azure Advisor's AI-driven recommendations can continue optimizing sizing on an ongoing basis (not just this one-time assessment), and Copilot for Azure can generate natural-language cost/capacity narratives for ongoing FinOps reviews - turning this static comparison into a continuously self-improving process.""",

    6: """FILESHARES BY HOST OS CATEGORY - ASSESSMENT

VM RETIREMENT IMPACT: 743 of 746 fileshares (99.6%) run on Windows hosts. Combined with the Fileshare Readiness slide's finding of 68 fileshare-only servers, this Windows-heavy, homogeneous footprint is an ideal Azure Files migration target - full VM decommission is achievable for any dedicated file server, eliminating Windows Server OS licensing and patch management for each one retired.

SECURITY RISK: A single dominant OS category (99.6% Windows) means any Windows-specific vulnerability affects nearly the entire fileshare estate simultaneously - this concentration risk is worth flagging even though these hosts weren't separately assessed for support status on this slide (cross-reference against the broader out-of-support findings elsewhere in the deck).

OTHER OBSERVATIONS: Because the OS mix is so homogeneous (743 Windows vs. 3 Linux/RHEL), the migration approach can be highly standardized - a single Azure Files migration playbook covers virtually the entire fileshare estate, with only 3 shares requiring special handling for the Linux/RHEL hosts.

AI OPPORTUNITY: Once consolidated onto Azure Files, this content becomes indexable by Azure AI Search and usable as a grounding source for Copilot-driven document search/summarization - unlocking discoverability across file content that today sits siloed on individual Windows file servers with no enterprise search capability.""",

    7: """DATABASE RESOURCES BY TYPE - ASSESSMENT

VM RETIREMENT IMPACT: 289 total database resources span SQL Server (252, 87.2%), MongoDB (19), MySQL (11), and PostgreSQL (7). Per the SQL Readiness and Non-SQL Database Readiness slides, up to 99 SQL-only servers and all 37 non-SQL database servers are candidates for full retirement via PaaS conversion (SQL MI, Cosmos DB, Flexible Server) - eliminating OS/DB engine licensing, patching, and backup/HA management for well over 130 servers combined.

SECURITY RISK: This slide is the aggregate view of a database estate where a meaningful share carries out-of-support exposure - see SQL Readiness (3 SQL 2014 instances, 186 more on an expiring support clock) and Non-SQL Database Readiness (14 of 37 instances, 38%, already out of support). Because these are databases (not just app servers), any unpatched vulnerability here risks direct exposure of the customer's actual data.

OTHER OBSERVATIONS: SQL Server's 87.2% concentration means it is by far the largest single lever in the database modernization story - licensing cost, support-status risk, and MI-readiness here dominate the overall business case math. The non-SQL engines are comparatively small in volume but should not be deprioritized given their disproportionately high out-of-support rate relative to their size.

AI OPPORTUNITY: Once consolidated onto Azure PaaS, this 289-instance database estate becomes the foundation for Microsoft Fabric integration, Copilot-assisted database operations (natural-language query generation, automated tuning recommendations), and RAG-based applications grounded directly against the customer's own transactional data - none of which is practical against a fragmented, partially out-of-support on-prem estate.""",

    8: """SQL READINESS - ASSESSMENT

SCALE AND PAAS READINESS: 252 SQL Server instances across 245 servers, with 244 instances (97%) targeted for Azure SQL Managed Instance (PaaS) and only 8 staying on SQL VM (IaaS). This is the strongest PaaS-readiness ratio of any workload in the assessment.

VM RETIREMENT IMPACT: 99 of the 245 SQL servers run nothing but SQL (SQL-only servers) - these are candidates for full retirement if their instances migrate to SQL MI. Per server eliminated: SQL Server Enterprise/Standard per-core licensing (often the single largest software cost line item in a customer's Microsoft estate) goes away entirely (with Azure Hybrid Benefit available to offset transition costs); DBA time spent on patching, backups, and HA/failover configuration is eliminated as SQL MI automates this natively; and VM compute/storage/OS licensing for all 99 boxes is retired. This 99-server figure is a floor - some of the remaining 146 mixed-role servers may also see reduced footprint.

SECURITY RISK - OUT OF SUPPORT: SQL Server 2014 (3 instances) has NO vendor security patching - any new CVE is a permanent, unpatchable exposure directly against customer data (not just service downtime). Treat as an immediate remediation priority regardless of broader timeline, despite the small count. Additionally, 186 instances (74% of the fleet) sit on Extended Support versions (2019: 111, 2016: 52, 2017: 23) - Extended Support still patches today but is a time-limited window, not a stable end-state; SQL Server 2016 in particular is at or near the end of its extended support lifecycle now. This is the largest looming risk in the entire assessment and should anchor the multi-year roadmap.

DATA QUALITY FLAG: 8 instances report version "17.0.1000.7" with Unknown support status - this doesn't match standard SQL Server release naming and should be investigated before finalizing wave assignments, since it could mask a hidden out-of-support instance.

OTHER OBSERVATIONS: 224 of 244 MI-targeted instances (92%) are "Ready w/ Conditions," not unconditionally Ready - get the specific blocking conditions (cross-database queries, linked servers, SQL Agent dependencies, unsupported features) from the assessment tool before committing a timeline. The 8 SQL VM-targeted instances won't retire their VMs, but still benefit from Azure-managed automated patching and Azure Hybrid Benefit licensing.

PRIORITIZATION: SQL 2014 first (zero-day-equivalent risk, small footprint), then the 99 SQL-only servers (fastest clean VM-retirement win), then the Extended Support fleet on a scheduled follow-up wave before support windows lapse, with the unknown version resolved in parallel.

AI OPPORTUNITY: Migrating to SQL MI unlocks Microsoft Fabric integration (native connectivity for analytics/Power BI/Copilot without complex ETL), Copilot in Azure SQL for natural-language query generation and automated performance tuning, and Azure AI-powered anomaly detection (Query Performance Insight, Intelligent Insights) that isn't available against on-prem, out-of-support SQL versions like 2014/2016. This turns 252 instances of operational data into a foundation for AI-driven database operations and downstream RAG/analytics use cases.""",

    9: """SQL COST AND LICENSING - ASSESSMENT

LICENSE OWNERSHIP: Customer owns 136 SQL Server Enterprise licenses (272 cores) and 652 SQL Server Standard licenses (1,304 cores) today via EA export - a substantial existing software investment that should inform the Azure Hybrid Benefit strategy during migration (reusing owned licenses against Azure compute reduces the effective cost delta below).

COST COMPARISON - CLEAR FINANCIAL CASE FOR SQL MI: 100% Lift & Shift (245 VMs, 0 SQL MI) costs $80.7K/mo ($968,154/yr). SQL MI (where ready) + Lift & Shift (8 VMs, 244 SQL MI) costs $156.7K/mo ($1,879,978/yr) - notably HIGHER in raw monthly spend. This is an important nuance: the MI-heavy path costs more in direct Azure consumption, but that additional spend buys full automation of patching, backup, and HA - meaning the true comparison must weigh this delta against the labor/licensing/downtime-risk savings from retiring 237 additional VMs (245 minus 8) and their associated DBA overhead, not just Azure bill line items. Recommend building a blended TCO model (Azure spend + retired license costs + reduced ops labor + risk-avoidance value) rather than comparing consumption cost alone.

VM RETIREMENT IMPACT: The SQL MI path retires all but 8 VMs (237 servers), each carrying its own licensing, patching burden, and DBA-hours - the true business case rests on offsetting the higher direct Azure MI cost against this dramatically reduced infrastructure/operational footprint.

RESOURCE EFFICIENCY: SQL MI path shows meaningfully lower resource consumption versus on-prem - cores (1,744 vs 1,721, roughly flat), memory (10.8 TB vs 16.7 TB, -35%), and storage (341.93 TB vs 1,069.50 TB, -68%) - reinforcing that the on-prem SQL estate is significantly overprovisioned relative to actual need.

SECURITY RISK CONTEXT: This cost analysis should not be evaluated in isolation from the SQL Readiness findings - 3 SQL 2014 instances (no patching) and 186 instances on an expiring Extended Support clock represent real breach/compliance exposure that a "stay on Lift & Shift only" decision does not resolve. The cost delta above should be weighed against the cost of a potential breach or audit finding tied to unpatched database infrastructure.

AI OPPORTUNITY: Beyond the direct cost comparison, migrating to SQL MI positions the customer to use Copilot in Azure SQL and Microsoft Fabric for cost/performance optimization recommendations on an ongoing basis, and Azure Cost Management + Copilot for Azure can continuously re-evaluate the Lift & Shift vs. SQL MI mix as workload patterns change - turning this static cost comparison into a living, AI-assisted FinOps process rather than a one-time snapshot.""",

    10: """NON-SQL DATABASE READINESS - ASSESSMENT

PAAS READINESS: All 37 instances (MongoDB: 19, MySQL: 11, PostgreSQL: 7) show 100% "Ready w/ Conditions" for PaaS migration (Cosmos DB, Flexible Server, Flexible Server respectively) - 0 Not Ready. This is the cleanest readiness profile in the entire assessment.

VM RETIREMENT IMPACT: Since instance-to-server ratio is essentially 1:1 across all three engines (19 MongoDB instances / 19 servers, 11 MySQL / 11 servers, 7 PostgreSQL / 7 servers), all 37 servers are candidates for full retirement if migrated. Per server: database engine licensing/support subscriptions, backup infrastructure, HA/replication tooling, and DBA patching/tuning time are all eliminated - shifted entirely to Microsoft's managed service responsibility. This is one of the cleanest, most quantifiable VM-elimination opportunities in the deck given the 1:1 ratio and zero "Not Ready" blockers.

SECURITY RISK - OUT OF SUPPORT: 14 of 37 instances (38%) are Out of Support today, plus 9 more (24%) Unknown - meaning up to 23 of 37 (62%) have no confirmed security posture. Out-of-support DATABASES carry materially higher risk than out-of-support app/web servers, because a breach here exposes the customer's actual data directly (exfiltration, corruption, ransomware targeting the data layer), not just service downtime. This should be framed as the headline driver of urgency - PaaS readiness explains how to fix it, but the out-of-support exposure explains why it cannot wait.

BY ENGINE:
- MySQL: 8 of 11 instances (73%) Out of Support - the highest risk density of any engine in this assessment. Smallest instance count of the three, meaning fastest time-to-risk-reduction. Recommend leading with MySQL.
- MongoDB: Largest footprint (19 servers), with 4 Out of Support and 5 Unknown (47% combined exposure). Migrating to Cosmos DB retires the most VMs of any engine here - the single largest cost/ops-savings line item in this modernization workstream.
- PostgreSQL: Smallest count (7 servers), but worst visibility gap - 3 of 7 (43%) Unknown support status. Recommend a targeted discovery pass before committing this engine to a migration wave, since "Unknown" could mean anything from current to critically exposed.

PRIORITIZATION: MySQL first (highest risk density + guaranteed retirement + smallest lift), MongoDB second (largest retirement/cost-savings volume), PostgreSQL third (resolve the Unknown-status gap via discovery before finalizing scope).

AI OPPORTUNITY: Migrating MongoDB, MySQL, and PostgreSQL workloads to their Azure PaaS equivalents (Cosmos DB, Flexible Server) enables native integration with Microsoft Fabric for unified analytics alongside the SQL Server estate, positions the data for Azure AI Search-grounded RAG applications, and unlocks AI-powered performance/anomaly monitoring that isn't available against on-prem, partially out-of-support database versions. This also closes a data-layer security gap that only widens the longer it's deferred.""",

    11: """WEBAPP READINESS - ASSESSMENT

SCALE: 870 total web applications (IIS: 866, 99.5%; Tomcat: 4, 0.5%) across 683 hosting servers, with 475 identified as dedicated WebApp servers (no other workload role).

VM RETIREMENT IMPACT: All 475 dedicated servers are candidates for full decommission if their apps migrate to Azure App Service (PaaS) - not resized, retired outright. Per server eliminated: Windows Server OS licensing goes away, patch cycles/backup jobs/monitoring configuration are eliminated, and the compute/storage/networking cost of that VM is replaced by App Service's consumption-based pricing (typically cheaper per-app than a dedicated VM). This should be quantified explicitly in the business case (VM compute + OS/software licensing + ops labor per server x 475) against the App Service Plan cost for the consolidated app set.

SECURITY RISK - OUT OF SUPPORT: 27 servers (confirmed Out of Support) + 84 servers (Unknown status) = 111 servers, 13% of the IIS fleet, represent an active business risk today, independent of migration timing. Out-of-support operating systems/runtimes receive no further security patches - any newly disclosed CVE against that OS/IIS version is a permanent, unpatchable exposure until the server is retired or upgraded. This creates breach/ransomware exposure and compliance audit findings (PCI, HIPAA, SOC2 all flag unsupported software as a control failure). The 84 "Unknown" servers are arguably the bigger concern since the true at-risk count could exceed 111. RECOMMENDATION: treat these 111 servers as migration wave one, framing the business case around risk elimination first, cost savings second - migrating them doesn't just delay the risk, it eliminates it permanently since the customer no longer patches the underlying OS at all.

CONTAINERIZATION PATH: 402 of 870 apps (46%) are flagged "Ready w/ Conditions" for AKS. AKS won't zero out VMs the way App Service does (nodes are still shared VMs), but still converts a 1:1 app-to-VM footprint into a many-to-few model, capturing most of the same licensing/ops-overhead savings at a smaller per-node scale. This also resolves the underlying OS support-risk problem for any at-risk servers that land in this bucket.

DATA GAP: Readiness assessment only covers 402 of 870 WebApps (46%) - the remaining ~468 apps are unassessed. Given that 111 servers are a known or suspected security risk today, closing this gap should happen BEFORE final wave planning, since at-risk servers with no confirmed migration path yet are the most urgent unknown in this deck.

OTHER OBSERVATIONS: ~208 servers are mixed-role (683 total minus 475 dedicated) and won't fully retire post-migration, but removing the webapp workload reduces their footprint and may simplify patching for any remaining out-of-support/unknown-status shared servers. Tomcat (4 apps, 2 servers, both Mainstream-supported) is a small, low-risk pilot candidate to validate the "migrate then decommission" process before the larger 475-server IIS wave.

PRIORITIZATION: Assess the remaining 468 unassessed apps first (especially any among the 111 at-risk servers), then prioritize the 111 out-of-support/unknown servers as wave one for risk elimination, then proceed with the broader 475-server App Service migration, using Tomcat as an early low-risk pilot.

AI OPPORTUNITY: Once these 870 web applications and their supporting infrastructure are modernized onto App Service/AKS, the customer can layer Copilot-assisted DevOps (GitHub Copilot for ongoing app maintenance, Azure Copilot for operational monitoring) and use Azure AI Search + Azure OpenAI to build RAG-based experiences grounded in application data and logs that are impractical to access safely on out-of-support, on-prem IIS servers today. This migration is the enabling step for bringing modern AI tooling to bear on a web application estate that currently sits outside the reach of cloud-native AI services.""",
}


def main():
    prs = Presentation(SRC)
    applied = 0
    for i, slide in enumerate(prs.slides, 1):
        note = NOTES.get(i)
        if not note:
            continue
        notes_slide = slide.notes_slide  # creates if not present
        notes_slide.notes_text_frame.text = note
        applied += 1
        print(f"Slide {i}: notes added ({len(note)} chars)")
    prs.save(OUT)
    print(f"\nSaved {OUT} with notes on {applied} slides.")


if __name__ == "__main__":
    main()
